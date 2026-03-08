"""文档解析器 - 支持多种格式的文档解析。

支持的格式：
- PDF: 使用 pymupdf4llm 解析
- DOCX: 使用 python-docx 解析
- PPTX: 使用 python-pptx 解析
- 图片: 使用视觉模型提取文字
- URL: 使用 Playwright + BeautifulSoup
- 文本文件: 直接读取

依赖安装：
    pip install pymupdf4llm python-docx python-pptx beautifulsoup4 lxml playwright
    playwright install chromium
"""

import logging
import os
import re
import tempfile
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """解析结果。"""
    filename: str
    file_type: str
    content: str
    metadata: dict
    success: bool
    error: Optional[str] = None


@dataclass
class URLParseResult:
    """URL 解析结果。"""
    url: str
    title: str
    content: str
    description: str
    metadata: dict


class DocumentParser:
    """文档解析器。"""

    # 支持的文件类型
    SUPPORTED_FILE_TYPES = {
        # 文档
        "pdf": "pdf",
        "docx": "docx",
        "doc": "docx",
        "pptx": "pptx",
        "ppt": "pptx",
        "txt": "text",
        "md": "markdown",
        "markdown": "markdown",
        "json": "json",
        "csv": "csv",
        "xlsx": "excel",
        "xls": "excel",
        # 图片
        "jpg": "image",
        "jpeg": "image",
        "png": "image",
        "gif": "image",
        "webp": "image",
        "bmp": "image",
        # 视频/音频
        "mp4": "video",
        "avi": "video",
        "mov": "video",
        "mp3": "audio",
        "wav": "audio",
        "m4a": "audio",
    }

    def __init__(self, vision_client=None, vision_model: str = "glm-4v-flash"):
        """初始化解析器。

        Args:
            vision_client: 视觉模型客户端（用于图片/视频处理）
            vision_model: 视觉模型名称
        """
        self.vision_client = vision_client
        self.vision_model = vision_model

    def parse(self, file_path: str) -> ParseResult:
        """解析文件。

        Args:
            file_path: 文件路径

        Returns:
            解析结果
        """
        path = Path(file_path)

        if not path.exists():
            return ParseResult(
                filename=path.name,
                file_type="unknown",
                content="",
                metadata={},
                success=False,
                error=f"文件不存在: {file_path}",
            )

        ext = path.suffix.lower().lstrip(".")
        file_type = self.SUPPORTED_FILE_TYPES.get(ext, "unknown")

        try:
            if file_type == "pdf":
                content = self._parse_pdf(file_path)
            elif file_type == "docx":
                content = self._parse_docx(file_path)
            elif file_type == "pptx":
                content = self._parse_pptx(file_path)
            elif file_type == "image":
                content = self._parse_image(file_path, path.name)
            elif file_type in ("video", "audio"):
                content = self._parse_media(file_path, path.name, file_type)
            elif file_type == "markdown":
                content = self._parse_markdown(file_path)
            elif file_type == "text":
                content = self._parse_text(file_path)
            elif file_type == "json":
                content = self._parse_json(file_path)
            elif file_type == "csv":
                content = self._parse_csv(file_path)
            elif file_type == "excel":
                content = self._parse_excel(file_path)
            else:
                content = self._parse_text(file_path)

            # 清理内容
            content = self._clean_text(content)

            return ParseResult(
                filename=path.name,
                file_type=file_type,
                content=content,
                metadata={
                    "file_size": path.stat().st_size,
                    "file_path": str(path.resolve()),
                },
                success=True,
            )

        except Exception as e:
            logger.error(f"❌ 解析失败 {file_path}: {e}")
            return ParseResult(
                filename=path.name,
                file_type=file_type,
                content="",
                metadata={},
                success=False,
                error=str(e),
            )

    def parse_url(self, url: str) -> URLParseResult:
        """解析 URL。

        Args:
            url: 网页 URL

        Returns:
            解析结果
        """
        return self._parse_url(url)

    # -------------------- 内部解析方法 --------------------

    def _parse_pdf(self, file_path: str) -> str:
        """解析 PDF 文件。"""
        try:
            # 尝试使用 pymupdf4llm - 新版本使用 to_markdown 函数
            import pymupdf4llm

            # 新版本 API: pymupdf4llm.to_markdown(doc_path)
            text = pymupdf4llm.to_markdown(file_path)

            # 如果结果为空或太短，回退到 pypdf
            if not text or len(text) < 100:
                logger.warning("pymupdf4llm 结果为空，尝试 pypdf")
                return self._parse_pdf_fallback(file_path)

            return text

        except ImportError:
            logger.warning("pymupdf4llm 未安装，使用 pypdf")
            return self._parse_pdf_fallback(file_path)
        except Exception as e:
            logger.warning(f"pymupdf4llm 解析失败: {e}，尝试 pypdf")
            return self._parse_pdf_fallback(file_path)

    def _parse_pdf_fallback(self, file_path: str) -> str:
        """使用 pypdf 解析 PDF（回退方案）。"""
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        text_parts = []

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- 第 {i + 1} 页 ---\n{page_text}")

        return "\n\n".join(text_parts)

    def _parse_docx(self, file_path: str) -> str:
        """解析 Word 文档。"""
        from docx import Document

        doc = Document(file_path)
        parts = []

        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                # 检查段落样式（标题）- 安全访问
                try:
                    style_name = para.style.name if para.style else ""
                    if style_name.startswith("Heading"):
                        parts.append(f"\n## {para.text}\n")
                    else:
                        parts.append(para.text)
                except Exception:
                    # 如果样式访问失败，直接添加文本
                    parts.append(para.text)

        # 提取表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                table_text.append(row_text)
            parts.append("\n--- 表格 ---\n")
            parts.append("\n".join(table_text))

        return "\n\n".join(parts)

    def _parse_pptx(self, file_path: str) -> str:
        """解析 PPT 文件。"""
        from pptx import Presentation

        prs = Presentation(file_path)
        parts = []

        for i, slide in enumerate(prs.slides):
            slide_text = [f"--- 第 {i + 1} 页 ---"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 检查是否为标题
                    if shape == slide.shapes.title:
                        slide_text.append(f"\n## {shape.text}")
                    else:
                        slide_text.append(shape.text)

            parts.append("\n".join(slide_text))

        return "\n\n".join(parts)

    def _parse_image(self, file_path: str, filename: str) -> str:
        """解析图片（使用视觉模型）。"""
        if self.vision_client is None:
            return f"[图片文件: {filename}]\n\n无法处理：未配置视觉模型客户端"

        try:
            import base64

            # 读取图片并编码为 base64
            with open(file_path, "rb") as f:
                image_data = base64.b64encode(f.read()).decode("utf-8")

            # 确定 MIME 类型
            ext = Path(file_path).suffix.lower()
            mime_types = {
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".gif": "image/gif",
                ".webp": "image/webp",
            }
            mime_type = mime_types.get(ext, "image/jpeg")

            # 调用视觉模型
            response = self.vision_client.chat.completions.create(
                model=self.vision_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{image_data}"
                                },
                            },
                            {
                                "type": "text",
                                "text": "请详细描述这张图片的内容。如果图片中有文字，请提取所有文字。请用中文回答，尽可能详细和准确。",
                            },
                        ],
                    }
                ],
            )

            extracted_text = response.choices[0].message.content
            return f"[图片文件: {filename}]\n\n{extracted_text}"

        except Exception as e:
            logger.error(f"❌ 图片处理失败: {e}")
            return f"[图片文件: {filename}]\n\n处理失败: {str(e)}"

    def _parse_media(self, file_path: str, filename: str, media_type: str) -> str:
        """解析视频/音频文件。"""
        return f"[{media_type} 文件: {filename}]\n\n视频/音频文件需要先提取文字。请使用 Whisper 等工具进行语音转文字处理。"

    def _parse_markdown(self, file_path: str) -> str:
        """解析 Markdown 文件。"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def _parse_text(self, file_path: str) -> str:
        """解析纯文本文件。"""
        # 尝试多种编码
        for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        # 最后尝试二进制读取
        with open(file_path, "rb") as f:
            return f.read().decode("utf-8", errors="ignore")

    def _parse_json(self, file_path: str) -> str:
        """解析 JSON 文件。"""
        import json

        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        return json.dumps(data, ensure_ascii=False, indent=2)

    def _parse_csv(self, file_path: str) -> str:
        """解析 CSV 文件。"""
        import pandas as pd

        df = pd.read_csv(file_path)
        return f"CSV 数据摘要:\n\n{df.to_string()}"

    def _parse_excel(self, file_path: str) -> str:
        """解析 Excel 文件。"""
        import pandas as pd

        # 读取所有 sheet
        sheets = pd.read_excel(file_path, sheet_name=None)
        parts = []

        for sheet_name, df in sheets.items():
            parts.append(f"--- Sheet: {sheet_name} ---")
            parts.append(df.to_string())
            parts.append("")

        return "\n\n".join(parts)

    def _parse_url(self, url: str) -> URLParseResult:
        """解析 URL 网页。"""
        try:
            # 尝试使用 playwright
            return self._parse_url_playwright(url)
        except Exception as e:
            logger.warning(f"Playwright 解析失败: {e}，尝试 requests")
            try:
                return self._parse_url_requests(url)
            except Exception as e2:
                logger.error(f"URL 解析完全失败: {e2}")
                return URLParseResult(
                    url=url,
                    title="",
                    content=f"网页解析失败: {str(e2)}",
                    description="",
                    metadata={},
                )

    def _parse_url_playwright(self, url: str) -> URLParseResult:
        """使用 Playwright 解析网页。"""
        from bs4 import BeautifulSoup

        # 动态导入 playwright
        from playwright.async_api import async_playwright

        async def extract_content():
            async with async_playwright() as p:
                browser = await p.chromium.launch()
                page = await browser.new_page()
                await page.goto(url, wait_until="networkidle")
                content = await page.content()
                await browser.close()
                return content

        # 使用 asyncio.run() 执行异步函数（兼容已运行的事件循环）
        import asyncio
        html_content = asyncio.run(extract_content())

        return self._extract_from_html(html_content, url)

    def _parse_url_requests(self, url: str) -> URLParseResult:
        """使用 requests 解析网页（回退方案）。"""
        import requests
        from bs4 import BeautifulSoup

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }

        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()

        return self._extract_from_html(response.text, url)

    def _extract_from_html(self, html: str, url: str) -> URLParseResult:
        """从 HTML 中提取内容。"""
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html, "lxml")

        # 提取标题
        title = ""
        if soup.title:
            title = soup.title.string or ""
        elif soup.find("h1"):
            title = soup.find("h1").get_text(strip=True)

        # 提取描述
        description = ""
        meta_desc = soup.find("meta", {"name": "description"})
        if meta_desc:
            description = meta_desc.get("content", "")

        # 提取正文（移除脚本和样式）
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()

        # 提取主要内容区域
        main_content = soup.find("main") or soup.find("article") or soup.body

        if main_content:
            # 获取文本，保留段落结构
            paragraphs = main_content.find_all("p")
            content_parts = []
            for p in paragraphs:
                text = p.get_text(strip=True)
                if text and len(text) > 20:  # 过滤短文本
                    content_parts.append(text)

            content = "\n\n".join(content_parts)
        else:
            content = soup.get_text(separator="\n", strip=True)

        return URLParseResult(
            url=url,
            title=title,
            content=content,
            description=description,
            metadata={
                "domain": urlparse(url).netloc,
            },
        )

    def _clean_text(self, text: str) -> str:
        """清理文本。"""
        if not text:
            return ""

        # 移除多余的空白
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)

        # 移除零宽字符
        text = re.sub(r'[\u200b-\u200f\u2028-\u202f]', '', text)

        return text.strip()


def create_parser(vision_client=None, vision_model: str = "glm-4v-flash") -> DocumentParser:
    """创建文档解析器。"""
    return DocumentParser(vision_client=vision_client, vision_model=vision_model)
