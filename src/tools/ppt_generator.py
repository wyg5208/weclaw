"""PPT 生成工具 — AI 生成演示文稿 PPT，支持多种风格模板。
"""

from __future__ import annotations

import logging
import shutil
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.tools.base import ActionDef, BaseTool, ToolResult, ToolResultStatus

logger = logging.getLogger(__name__)


# 风格配色方案
STYLE_COLORS = {
    "business": {
        "primary": RGBColor(0x00, 0x52, 0x8A),
        "secondary": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x00, 0x96, 0xD6),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "academic": {
        "primary": RGBColor(0x1A, 0x23, 0x7E),
        "secondary": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0xC6, 0x28, 0x28),
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
    },
    "creative": {
        "primary": RGBColor(0xE9, 0x1E, 0x63),
        "secondary": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0xFF, 0xC1, 0x07),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "minimal": {
        "primary": RGBColor(0x21, 0x21, 0x21),
        "secondary": RGBColor(0x75, 0x75, 0x75),
        "accent": RGBColor(0x42, 0x42, 0x42),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
}

# 增强配色方案 - 更丰富的渐变色彩
STYLE_COLORS_ENHANCED = {
    "gradient_blue": {
        "primary": RGBColor(0x00, 0x52, 0x8A),
        "secondary": RGBColor(0x1E, 0x88, 0xE5),
        "accent": RGBColor(0x00, 0x96, 0xD6),
        "highlight": RGBColor(0x4F, 0xC3, 0xF7),
        "bg": RGBColor(0xF0, 0xF8, 0xFF),
    },
    "gradient_purple": {
        "primary": RGBColor(0x6A, 0x1B, 0x9A),
        "secondary": RGBColor(0x7B, 0x1F, 0xA2),
        "accent": RGBColor(0xE1, 0xB5, 0xE7),
        "highlight": RGBColor(0xCE, 0x93, 0xD8),
        "bg": RGBColor(0xF3, 0xE5, 0xF5),
    },
    "gradient_orange": {
        "primary": RGBColor(0xE6, 0x5C, 0x00),
        "secondary": RGBColor(0xF5, 0x7C, 0x00),
        "accent": RGBColor(0xFF, 0xB7, 0x4D),
        "highlight": RGBColor(0xFF, 0xE0, 0x82),
        "bg": RGBColor(0xFF, 0xF8, 0xF0),
    },
    "gradient_green": {
        "primary": RGBColor(0x2E, 0x7D, 0x32),
        "secondary": RGBColor(0x43, 0xA0, 0x47),
        "accent": RGBColor(0x81, 0xC7, 0x84),
        "highlight": RGBColor(0xA5, 0xD6, 0xA7),
        "bg": RGBColor(0xF1, 0xF8, 0xE9),
    },
    "gradient_red": {
        "primary": RGBColor(0xC6, 0x28, 0x28),
        "secondary": RGBColor(0xE5, 0x39, 0x35),
        "accent": RGBColor(0xEF, 0x53, 0x50),
        "highlight": RGBColor(0xFF, 0x79, 0x6C),
        "bg": RGBColor(0xFF, 0xEB, 0xEE),
    },
}

# 图标映射 - 用于内容页装饰
ICONS_MAP = {
    "star": "★",
    "check": "✓",
    "arrow": "➤",
    "bullet": "●",
    "diamond": "◆",
    "fire": "🔥",
    "rocket": "🚀",
    "light": "💡",
    "target": "🎯",
    "chart": "📊",
    "money": "💰",
    "globe": "🌐",
    "tech": "⚙",
    "people": "👥",
    "idea": "💡",
    "warning": "⚠",
    "info": "ℹ",
}


class PPTTool(BaseTool):
    """PPT 生成工具。

    支持 AI 生成演示文稿 PPT，包含多种风格模板。
    - generate_ppt: 根据主题和大纲生成完整 PPT（简单版）
    - generate_rich_ppt: 生成图文并茂的 PPT（增强版）
    - add_slide: 向已有 PPT 添加幻灯片
    - export_pdf: 将 PPT 导出为 PDF
    """

    name = "ppt_generator"
    emoji = "📊"
    title = "PPT生成"
    description = "AI 生成演示文稿 PPT，支持多种风格模板、图文并茂"
    timeout = 120

    def __init__(self, output_dir: str = "") -> None:
        """初始化 PPT 生成工具。

        Args:
            output_dir: 输出目录，默认为项目的 generated/日期/ 目录
        """
        super().__init__()
        self.output_dir = (
            Path(output_dir)
            if output_dir
            else Path(__file__).parent.parent.parent / "generated" / datetime.now().strftime("%Y-%m-%d")
        )
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def get_actions(self) -> list[ActionDef]:
        return [
            ActionDef(
                name="generate_rich_ppt",
                description="【推荐】生成图文并茂的专业PPT，支持丰富布局、图标卡片、图片网格、图表表格",
                parameters={
                    "topic": {
                        "type": "string",
                        "description": "PPT主题，如'人工智能发展趋势报告'",
                    },
                    "slides_config": {
                        "type": "string",
                        "description": """幻灯片配置，JSON格式数组。每项格式：
{
    "type": "title|section|content|chart|table|icon_cards|grid_images|image|thank",
    "title": "页面标题",
    "subtitle": "副标题（可选）",
    "content": ["要点1", "要点2", "要点3"],  // content类型
    "layout": "single|double|left_img|right_img|grid_2|grid_3|grid_2x2|grid_3x2|cards_3|cards_4",  // 内容布局
    "icons": ["star", "rocket", "chart"],  // 图标列表（可选）
    "cards": [{"icon": "star", "title": "标题", "description": "描述"}],  // icon_cards类型
    "image_count": 4,  // grid_images类型的图片数量
    "images": ["path1.jpg", "path2.png"],  // 图片路径列表（可选）
    "image_path": "path/to/image.jpg",  // image类型的单张图片路径
    "image_layout": "left_half|right_half|left_large|right_large|full_bg|top_bottom",  // image类型布局
    "image_prompt": "配图描述（可选）",  // 用于生成配图
    "chart": {"type": "bar|pie|line", "labels": [], "values": []},  // chart类型
    "table": {"headers": [], "rows": [[]]}  // table类型
}""",
                    },
                    "style": {
                        "type": "string",
                        "description": "配色风格: business/academic/creative/minimal",
                        "enum": ["business", "academic", "creative", "minimal"],
                        "default": "business",
                    },
                    "include_images": {
                        "type": "boolean",
                        "description": "是否生成配图（需配置图片生成工具）",
                        "default": False,
                    },
                    "output_filename": {
                        "type": "string",
                        "description": "输出文件名（不含扩展名）",
                    },
                },
                required_params=["topic", "slides_config"],
            ),
            ActionDef(
                name="generate_ppt",
                description="AI 生成完整 PPT 演示文稿，根据主题和大纲自动创建多页幻灯片",
                parameters={
                    "topic": {
                        "type": "string",
                        "description": "PPT 主题",
                    },
                    "outline": {
                        "type": "string",
                        "description": "大纲内容，每行一个要点",
                    },
                    "style": {
                        "type": "string",
                        "description": "风格: business/academic/creative/minimal",
                        "enum": ["business", "academic", "creative", "minimal"],
                    },
                    "slide_count": {
                        "type": "integer",
                        "description": "幻灯片数量，默认8",
                    },
                    "output_filename": {
                        "type": "string",
                        "description": "输出文件名（不含扩展名）",
                    },
                },
                required_params=["topic"],
            ),
            ActionDef(
                name="add_slide",
                description="向已有 PPT 文件添加新的幻灯片",
                parameters={
                    "ppt_path": {
                        "type": "string",
                        "description": "已有 PPT 文件路径",
                    },
                    "title": {
                        "type": "string",
                        "description": "幻灯片标题",
                    },
                    "content": {
                        "type": "string",
                        "description": "幻灯片内容，支持多行，每行一个要点",
                    },
                    "layout": {
                        "type": "string",
                        "description": "布局: title/content/two_column/blank",
                        "enum": ["title", "content", "two_column", "blank"],
                    },
                },
                required_params=["ppt_path", "title", "content"],
            ),
            ActionDef(
                name="export_pdf",
                description="将 PPT 文件导出为 PDF 格式",
                parameters={
                    "ppt_path": {
                        "type": "string",
                        "description": "PPT 文件路径",
                    },
                    "output_filename": {
                        "type": "string",
                        "description": "输出 PDF 文件名（不含扩展名）",
                    },
                },
                required_params=["ppt_path"],
            ),
            ActionDef(
                name="add_chart_slide",
                description="添加图表幻灯片（柱状图、饼图、折线图）",
                parameters={
                    "ppt_path": {
                        "type": "string",
                        "description": "已有 PPT 文件路径",
                    },
                    "title": {
                        "type": "string",
                        "description": "幻灯片标题",
                    },
                    "chart_type": {
                        "type": "string",
                        "description": "图表类型",
                        "enum": ["bar", "pie", "line"],
                    },
                    "chart_data": {
                        "type": "string",
                        "description": "图表数据，JSON格式。如：{\"labels\":[\"A\",\"B\",\"C\"],\"values\":[10,20,30]}",
                    },
                    "style": {
                        "type": "string",
                        "description": "配色风格: business/academic/creative/minimal",
                        "enum": ["business", "academic", "creative", "minimal"],
                    },
                },
                required_params=["ppt_path", "title", "chart_type", "chart_data"],
            ),
            ActionDef(
                name="add_table_slide",
                description="添加表格幻灯片",
                parameters={
                    "ppt_path": {
                        "type": "string",
                        "description": "已有 PPT 文件路径",
                    },
                    "title": {
                        "type": "string",
                        "description": "幻灯片标题",
                    },
                    "table_data": {
                        "type": "string",
                        "description": "表格数据，JSON格式。格式：[{\"headers\":[\"列1\",\"列2\"]},{\"row1\":[\"数据1\",\"数据2\"]}]",
                    },
                    "style": {
                        "type": "string",
                        "description": "配色风格: business/academic/creative/minimal",
                        "enum": ["business", "academic", "creative", "minimal"],
                    },
                },
                required_params=["ppt_path", "title", "table_data"],
            ),
            ActionDef(
                name="add_section_divider",
                description="添加章节分隔页",
                parameters={
                    "ppt_path": {
                        "type": "string",
                        "description": "已有 PPT 文件路径",
                    },
                    "section_title": {
                        "type": "string",
                        "description": "章节标题",
                    },
                    "section_number": {
                        "type": "integer",
                        "description": "章节编号",
                    },
                    "style": {
                        "type": "string",
                        "description": "配色风格: business/academic/creative/minimal",
                        "enum": ["business", "academic", "creative", "minimal"],
                    },
                },
                required_params=["ppt_path", "section_title"],
            ),
        ]

    async def execute(self, action: str, params: dict[str, Any]) -> ToolResult:
        if action == "generate_rich_ppt":
            return await self._generate_rich_ppt(params)
        elif action == "generate_ppt":
            return self._generate_ppt(params)
        elif action == "add_slide":
            return self._add_slide(params)
        elif action == "export_pdf":
            return self._export_pdf(params)
        elif action == "add_chart_slide":
            return self._add_chart_slide(params)
        elif action == "add_table_slide":
            return self._add_table_slide(params)
        elif action == "add_section_divider":
            return self._add_section_divider(params)
        else:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"不支持的动作: {action}",
            )

    def _get_style_colors(self, style: str) -> dict[str, RGBColor]:
        """获取指定风格的配色方案。"""
        return STYLE_COLORS.get(style, STYLE_COLORS["business"])

    def _set_slide_background(self, slide, bg_color: RGBColor) -> None:
        """设置幻灯片背景色。"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def _add_title_slide(self, prs: Presentation, topic: str, colors: dict[str, RGBColor]) -> None:
        """添加标题页。"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors["bg"])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = topic
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # 副标题/日期
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(20)
        p.font.color.rgb = colors["secondary"]
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_title: str,
        subtitle: str,
        colors: dict[str, RGBColor],
        bullet_points: list[str] | None = None,
    ) -> None:
        """添加内容页。"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors["bg"])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = colors["accent"]
        p.alignment = PP_ALIGN.LEFT

        # 内容区域（要点列表）
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        if bullet_points:
            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.font.color.rgb = colors["secondary"]
                p.space_after = Pt(12)
        else:
            # 默认占位内容
            p = content_frame.paragraphs[0]
            p.text = "• 请在此处添加内容"
            p.font.size = Pt(20)
            p.font.color.rgb = colors["secondary"]

    def _add_thank_slide(self, prs: Presentation, colors: dict[str, RGBColor]) -> None:
        """添加感谢页。"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors["bg"])

        # 感谢文字
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "感谢聆听"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # 副文字
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = "THANK YOU"
        p.font.size = Pt(24)
        p.font.color.rgb = colors["accent"]
        p.alignment = PP_ALIGN.CENTER

    async def _generate_rich_ppt(self, params: dict[str, Any]) -> ToolResult:
        """生成图文并茂的 PPT（增强版）。

        支持：
        - 结构化幻灯片配置
        - 多种布局类型（单栏/双栏/左图右文/右图左文）
        - 图表插入（柱状图/饼图/折线图）
        - 表格插入
        - 可选配图生成
        """
        import json
        import logging

        topic = params.get("topic", "").strip()
        if not topic:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT 主题不能为空",
            )

        slides_config_str = params.get("slides_config", "").strip()
        if not slides_config_str:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="slides_config 不能为空，请提供 JSON 格式的幻灯片配置",
            )

        style = params.get("style", "business")
        include_images = params.get("include_images", False)
        output_filename = params.get("output_filename", "").strip()

        # 解析幻灯片配置
        try:
            slides_config = json.loads(slides_config_str)
        except json.JSONDecodeError as e:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"slides_config JSON 解析失败: {e}",
            )

        if not isinstance(slides_config, list):
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="slides_config 必须是 JSON 数组格式",
            )

        try:
            # 初始化 PPT
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9
            prs.slide_height = Inches(7.5)

            colors = self._get_style_colors(style)

            # 逐页生成幻灯片
            for i, slide_config in enumerate(slides_config):
                slide_type = slide_config.get("type", "content")
                layout_type = slide_config.get("layout", "single")

                # 根据类型选择布局
                layout_idx = self._get_rich_layout(slide_type, layout_type)
                slide_layout = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(slide_layout)
                self._set_slide_background(slide, colors["bg"])

                # 清除继承的占位符默认文本
                self._clear_placeholders(slide)

                # 根据类型填充内容
                if slide_type == "title":
                    self._fill_rich_title_slide(slide, slide_config, colors)
                elif slide_type == "section":
                    self._fill_rich_section_slide(slide, slide_config, colors)
                elif slide_type == "content":
                    self._fill_rich_content_slide(slide, slide_config, colors, include_images)
                elif slide_type == "chart":
                    self._fill_rich_chart_slide(slide, slide_config, colors)
                elif slide_type == "table":
                    self._fill_rich_table_slide(slide, slide_config, colors)
                elif slide_type == "icon_cards":
                    self._fill_rich_icon_cards_slide(slide, slide_config, colors)
                elif slide_type == "grid_images":
                    self._fill_rich_grid_images_slide(slide, slide_config, colors)
                elif slide_type == "image":
                    self._fill_image_slide(slide, slide_config, colors)
                elif slide_type == "thank":
                    self._fill_rich_thank_slide(slide, slide_config, colors)
                else:
                    # 默认作为内容页处理
                    self._fill_rich_content_slide(slide, slide_config, colors, include_images)

            # 生成文件名
            if not output_filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                safe_topic = "".join(c for c in topic[:20] if c.isalnum() or c in " _-")
                output_filename = f"ppt_rich_{safe_topic}_{timestamp}"

            output_path = self.output_dir / f"{output_filename}.pptx"
            prs.save(str(output_path))

            file_size = output_path.stat().st_size
            slide_count = len(prs.slides)

            logger.info("图文并茂 PPT 生成成功: %s (%d 页)", output_path, slide_count)

            return ToolResult(
                status=ToolResultStatus.SUCCESS,
                output=f"✅ 图文并茂的 PPT 生成完成\n📁 文件: {output_path.name}\n📊 大小: {file_size} 字节\n📄 共 {slide_count} 页\n🎨 风格: {style}",
                data={
                    "file_path": str(output_path),
                    "file_name": output_path.name,
                    "file_size": file_size,
                    "slide_count": slide_count,
                    "style": style,
                },
            )

        except Exception as e:
            logger.error("图文并茂 PPT 生成失败: %s", e, exc_info=True)
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT 生成失败: {e}",
            )

    def _get_rich_layout(self, slide_type: str, layout_type: str) -> int:
        """根据幻灯片类型和布局类型获取布局索引。"""
        # 布局映射 - 使用精确匹配和默认值
        type_defaults = {
            "title": 0,
            "section": 2,
            "content": 1,
            "chart": 5,
            "table": 5,
            "thank": 0,
            "blank": 6,
        }

        # 类型 + 布局组合映射
        combined_map = {
            ("title", "single"): 0,
            ("title", "double"): 0,
            ("title", "left_img"): 0,
            ("title", "right_img"): 0,
            ("section", "single"): 2,
            ("section", "double"): 2,
            ("section", "left_img"): 2,
            ("section", "right_img"): 2,
            ("content", "single"): 1,
            ("content", "double"): 3,
            ("content", "left_img"): 7,
            ("content", "right_img"): 8,
            ("chart", "single"): 5,
            ("chart", "double"): 5,
            ("chart", "left_img"): 5,
            ("chart", "right_img"): 5,
            ("table", "single"): 5,
            ("table", "double"): 5,
            ("table", "left_img"): 5,
            ("table", "right_img"): 5,
            ("thank", "single"): 0,
            ("thank", "double"): 0,
            ("thank", "left_img"): 0,
            ("thank", "right_img"): 0,
        }

        # 优先查找精确匹配
        key = (slide_type, layout_type)
        if key in combined_map:
            return combined_map[key]

        # 其次按类型查找默认布局
        if slide_type in type_defaults:
            return type_defaults[slide_type]

        # 默认空白布局
        return 6

    def _clear_placeholders(self, slide) -> None:
        """清除幻灯片中所有继承的占位符默认文本。

        Args:
            slide: 幻灯片对象
        """
        for shape in slide.shapes:
            if shape.is_placeholder:
                # 获取占位符的文本框架并清除内容
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ""

    def _fill_rich_title_slide(self, slide, config: dict, colors: dict) -> None:
        """填充标题页。"""
        title = config.get("title", "")
        subtitle = config.get("subtitle", "")

        # 使用占位符
        if hasattr(slide, 'shapes') and slide.shapes.title:
            slide.shapes.title.text = title

        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle or datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(20)
        p.font.color.rgb = colors["secondary"]
        p.alignment = PP_ALIGN.CENTER

    def _fill_rich_section_slide(self, slide, config: dict, colors: dict) -> None:
        """填充章节分隔页。"""
        title = config.get("title", "")
        subtitle = config.get("subtitle", "")
        section_num = config.get("section_number", "")

        # 使用占位符或文本框
        if hasattr(slide, 'shapes') and slide.shapes.title:
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 设置背景色
        self._set_slide_background(slide, colors["primary"])

        # 章节编号
        if section_num:
            num_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.8), Inches(12.333), Inches(0.8)
            )
            num_frame = num_box.text_frame
            p = num_frame.paragraphs[0]
            p.text = f"第 {section_num} 部分"
            p.font.size = Pt(24)
            p.font.color.rgb = colors["accent"]
            p.alignment = PP_ALIGN.CENTER

    def _fill_rich_content_slide(self, slide, config: dict, colors: dict, include_images: bool = False) -> None:
        """填充内容页。"""
        title = config.get("title", "")
        subtitle = config.get("subtitle", "")
        content = config.get("content", [])
        layout = config.get("layout", "single")
        icons = config.get("icons", ["star", "chart", "rocket"])

        # 顶部装饰
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors["primary"]
        top_bar.line.fill.background()

        # 标题图标
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = ICONS_MAP.get(icons[0], "★")
        p.font.size = Pt(28)
        p.font.color.rgb = colors["accent"]

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.0), Inches(11), Inches(0.5))
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = colors["accent"]

        # 分隔线
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = colors.get("highlight", colors["accent"])
        line.line.fill.background()

        # 内容区域
        if layout == "double":
            # 双栏布局
            self._fill_two_column_content_with_icons(slide, content, icons, colors)
        else:
            # 单栏布局
            self._fill_single_column_content_with_icons(slide, content, icons, colors)

    def _fill_single_column_content(self, slide, content: list, colors: dict) -> None:
        """填充单栏内容。"""
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.2), Inches(11.733), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        if content:
            for i, point in enumerate(content):
                p = content_frame.paragraphs[i] if i == 0 else content_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.font.color.rgb = colors["secondary"]
                p.space_after = Pt(12)
        else:
            p = content_frame.paragraphs[0]
            p.text = "• 请在此处添加内容"
            p.font.size = Pt(20)
            p.font.color.rgb = colors["secondary"]

    def _fill_single_column_content_with_icons(self, slide, content: list, icons: list, colors: dict) -> None:
        """填充单栏内容（带图标装饰）。"""
        y_start = 1.8
        for i, point in enumerate(content):
            y = y_start + i * 1.0

            # 圆点背景
            dot = slide.shapes.add_shape(9, Inches(0.5), Inches(y + 0.1), Inches(0.35), Inches(0.35))
            dot.fill.solid()
            dot.fill.fore_color.rgb = colors["primary"]
            dot.line.fill.background()

            # 图标
            icon_char = ICONS_MAP.get(icons[i] if i < len(icons) else "check", "✓")
            icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.05), Inches(0.35), Inches(0.4))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon_char
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            # 内容
            content_box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11.5), Inches(0.9))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = colors["secondary"]

    def _fill_two_column_content(self, slide, content: list, colors: dict) -> None:
        """填充双栏内容。"""
        mid = len(content) // 2
        left_content = content[:mid] if mid > 0 else content
        right_content = content[mid:] if mid > 0 else []

        # 左栏
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.0), Inches(5.9), Inches(5.0)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, point in enumerate(left_content):
            p = left_frame.paragraphs[i] if i == 0 else left_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = colors["secondary"]
            p.space_after = Pt(10)

        # 右栏
        right_box = slide.shapes.add_textbox(
            Inches(6.9), Inches(2.0), Inches(5.9), Inches(5.0)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, point in enumerate(right_content):
            p = right_frame.paragraphs[i] if i == 0 else right_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = colors["secondary"]
            p.space_after = Pt(10)

    def _fill_two_column_content_with_icons(self, slide, content: list, icons: list, colors: dict) -> None:
        """填充双栏内容（带图标装饰）。"""
        mid = len(content) // 2
        left_content = content[:mid] if mid > 0 else content
        right_content = content[mid:] if mid > 0 else []

        # 左栏
        y_start = 1.8
        for i, point in enumerate(left_content):
            y = y_start + i * 1.2

            # 圆点
            dot = slide.shapes.add_shape(9, Inches(0.5), Inches(y + 0.1), Inches(0.3), Inches(0.3))
            dot.fill.solid()
            dot.fill.fore_color.rgb = colors["primary"]
            dot.line.fill.background()

            # 图标
            icon_char = ICONS_MAP.get(icons[i] if i < len(icons) else "check", "✓")
            icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.05), Inches(0.3), Inches(0.35))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon_char
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            # 内容
            content_box = slide.shapes.add_textbox(Inches(0.95), Inches(y), Inches(5.4), Inches(1.0))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = colors["secondary"]

        # 右栏
        for i, point in enumerate(right_content):
            idx = i + mid
            y = y_start + i * 1.2

            # 圆点
            dot = slide.shapes.add_shape(9, Inches(6.8), Inches(y + 0.1), Inches(0.3), Inches(0.3))
            dot.fill.solid()
            dot.fill.fore_color.rgb = colors["primary"]
            dot.line.fill.background()

            # 图标
            icon_char = ICONS_MAP.get(icons[idx] if idx < len(icons) else "check", "✓")
            icon_box = slide.shapes.add_textbox(Inches(6.8), Inches(y + 0.05), Inches(0.3), Inches(0.35))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon_char
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            # 内容
            content_box = slide.shapes.add_textbox(Inches(7.25), Inches(y), Inches(5.4), Inches(1.0))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = colors["secondary"]

    def _fill_rich_chart_slide(self, slide, config: dict, colors: dict) -> None:
        """填充图表页。"""
        import json

        title = config.get("title", "数据分析")
        chart_data = config.get("chart", {})

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        # 创建图表
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE

            chart_type_str = chart_data.get("type", "bar")
            chart_type_map = {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "pie": XL_CHART_TYPE.PIE,
                "line": XL_CHART_TYPE.LINE,
            }
            chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.BAR_CLUSTERED)

            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get("labels", [])
            chart_data_obj.add_series("数据", chart_data.get("values", []))

            x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(10), Inches(5)
            slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data_obj)

        except Exception as e:
            logger.warning(f"图表创建失败: {e}，使用占位文本")
            # 降级为文本占位
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.0), Inches(12), Inches(4)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            p = content_frame.paragraphs[0]
            p.text = f"📊 图表: {chart_type_str}\n标签: {chart_data.get('labels', [])}\n数值: {chart_data.get('values', [])}"
            p.font.size = Pt(16)
            p.font.color.rgb = colors["secondary"]

    def _fill_rich_table_slide(self, slide, config: dict, colors: dict) -> None:
        """填充表格页。"""
        import json

        title = config.get("title", "数据表格")
        table_data = config.get("table", {})

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        # 创建表格
        try:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            if headers and rows:
                num_rows = len(rows) + 1
                num_cols = max(len(headers), max(len(row) for row in rows) if rows else 0)

                x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12), Inches(5)
                table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy)
                table = table_shape.table

                # 设置列宽
                col_width = int(cx / num_cols)
                for col_idx in range(num_cols):
                    table.columns[col_idx].width = col_width

                # 填充表头
                for col_idx, header in enumerate(headers[:num_cols]):
                    cell = table.cell(0, col_idx)
                    cell.text = str(header)
                    for para in cell.text_frame.paragraphs:
                        para.font.bold = True
                        para.font.size = Pt(12)
                        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = colors["primary"]

                # 填充数据
                for row_idx, row in enumerate(rows):
                    for col_idx, value in enumerate(row[:num_cols]):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(value)
                        for para in cell.text_frame.paragraphs:
                            para.font.size = Pt(11)
                            para.font.color.rgb = colors["secondary"]
            else:
                logger.warning("表格数据为空")

        except Exception as e:
            logger.warning(f"表格创建失败: {e}，使用占位文本")

    def _fill_rich_thank_slide(self, slide, config: dict, colors: dict) -> None:
        """填充感谢页。"""
        title = config.get("title", "感谢聆听")
        subtitle = config.get("subtitle", "THANK YOU")

        # 感谢文字
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # 副文字
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = colors["accent"]
        p.alignment = PP_ALIGN.CENTER

    def _fill_rich_icon_cards_slide(self, slide, config: dict, colors: dict) -> None:
        """填充图标卡片页。"""
        title = config.get("title", "核心要点")
        cards = config.get("cards", [])
        icons = config.get("icons", ["star", "chart", "rocket", "light"])

        # 顶部装饰
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors["primary"]
        top_bar.line.fill.background()

        # 标题图标
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = "★"
        p.font.size = Pt(28)
        p.font.color.rgb = colors["accent"]

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        # 分隔线
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = colors.get("highlight", colors["accent"])
        line.line.fill.background()

        # 计算卡片布局
        card_count = len(cards) if cards else 3
        if card_count == 3:
            w, h = 3.9, 5.5
            gap = 0.3
            positions = [(0.5, 1.4), (0.5 + w + gap, 1.4), (0.5 + (w + gap) * 2, 1.4)]
        elif card_count == 4:
            w, h = 5.8, 2.5
            gap = 0.3
            positions = [(0.5, 1.4), (0.5 + w + gap, 1.4), (0.5, 1.4 + h + gap), (0.5 + w + gap, 1.4 + h + gap)]
        elif card_count == 6:
            w, h = 3.9, 2.5
            gap = 0.2
            positions = [
                (0.5, 1.4), (0.5 + w + gap, 1.4), (0.5 + (w + gap) * 2, 1.4),
                (0.5, 1.4 + h + gap), (0.5 + w + gap, 1.4 + h + gap), (0.5 + (w + gap) * 2, 1.4 + h + gap)
            ]
        else:
            w, h = 5.8, 5.5
            positions = [(0.5, 1.4)]

        # 绘制卡片
        for i in range(min(card_count, len(positions))):
            x, y = positions[i]
            if cards and i < len(cards):
                card = cards[i]
                icon = card.get("icon", icons[i % len(icons)] if icons else "star")
                card_title = card.get("title", f"要点 {i+1}")
                card_desc = card.get("description", "")
            else:
                icon = icons[i % len(icons)] if icons else "star"
                card_title = f"要点 {i+1}"
                card_desc = f"这是第 {i+1} 个要点的详细描述"

            # 卡片背景
            shape = slide.shapes.add_shape(12, Inches(x), Inches(y), Inches(w), Inches(h))  # 12 = ROUNDED_RECTANGLE
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shape.line.color.rgb = colors["primary"]
            shape.line.width = Pt(1.5)

            # 顶部装饰条
            top = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.15))
            top.fill.solid()
            top.fill.fore_color.rgb = colors["primary"]
            top.line.fill.background()

            # 图标
            icon_char = ICONS_MAP.get(icon, icon)
            icon_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.4), Inches(w), Inches(1))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon_char
            p.font.size = Pt(48)
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = colors["accent"]

            # 标题
            title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.5), Inches(w - 0.4), Inches(0.6))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = colors["primary"]

            # 描述
            if card_desc:
                desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.1), Inches(w - 0.4), Inches(h - 2.5))
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = card_desc
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = colors["secondary"]

    def _fill_rich_grid_images_slide(self, slide, config: dict, colors: dict) -> None:
        """填充图片网格页（支持真实图片）。"""
        title = config.get("title", "配图展示")
        image_count = config.get("image_count", 4)
        images = config.get("images", [])  # 图片路径列表

        # 顶部装饰
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors["primary"]
        top_bar.line.fill.background()

        # 标题图标
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = "🖼️"
        p.font.size = Pt(28)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        # 图片网格布局
        y_start = 1.3
        height = 5.5
        highlight_color = colors.get("highlight", RGBColor(0xE3, 0xF2, 0xFD))

        def add_image_placeholder(x, y, w, h, label, image_path=None):
            # 如果有真实图片路径，插入图片
            if image_path and Path(image_path).exists():
                try:
                    # 计算保持纵横比的尺寸
                    slide.shapes.add_picture(
                        image_path,
                        Inches(x), Inches(y), Inches(w), Inches(h)
                    )
                    return
                except Exception as e:
                    logger.warning(f"插入图片失败: {e}")
            
            # 阴影
            shadow = slide.shapes.add_shape(1, Inches(x + 0.05), Inches(y + 0.05), Inches(w), Inches(h))
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = colors["secondary"]
            shadow.line.fill.background()

            # 主框
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = highlight_color
            shape.line.color.rgb = colors["primary"]
            shape.line.width = Pt(2)

            # 中心图标
            icon_box = slide.shapes.add_textbox(Inches(x), Inches(y + h/2 - 0.6), Inches(w), Inches(1))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = "📷"
            p.font.size = Pt(48)
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = colors["accent"]

            # 标签
            label_box = slide.shapes.add_textbox(Inches(x), Inches(y + h/2 + 0.4), Inches(w), Inches(0.5))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = colors["primary"]

        if image_count == 1:
            add_image_placeholder(0.5, y_start, 12.333, height, "核心配图", images[0] if len(images) > 0 else None)
        elif image_count == 2:
            add_image_placeholder(0.5, y_start, 6.0, height, "配图 1", images[0] if len(images) > 0 else None)
            add_image_placeholder(6.8, y_start, 6.0, height, "配图 2", images[1] if len(images) > 1 else None)
        elif image_count == 3:
            add_image_placeholder(0.5, y_start, 6.0, height, "配图 1", images[0] if len(images) > 0 else None)
            add_image_placeholder(6.8, y_start, 5.8, height/2 - 0.1, "配图 2", images[1] if len(images) > 1 else None)
            add_image_placeholder(6.8, y_start + height/2 + 0.1, 5.8, height/2 - 0.1, "配图 3", images[2] if len(images) > 2 else None)
        elif image_count == 4:
            w = 6.0
            h = height/2 - 0.1
            gap = 0.2
            positions = [(0.5, y_start), (0.5 + w + gap, y_start), (0.5, y_start + h + gap), (0.5 + w + gap, y_start + h + gap)]
            for i, (x, y) in enumerate(positions):
                add_image_placeholder(x, y, w, h, f"配图 {i+1}", images[i] if i < len(images) else None)
        elif image_count == 6:
            w = 3.9
            h = 2.5
            gap = 0.2
            positions = [
                (0.5, y_start), (0.5 + w + gap, y_start), (0.5 + (w + gap) * 2, y_start),
                (0.5, y_start + h + gap), (0.5 + w + gap, y_start + h + gap), (0.5 + (w + gap) * 2, y_start + h + gap)
            ]
            for i, (x, y) in enumerate(positions):
                add_image_placeholder(x, y, w, h, f"配图 {i+1}", images[i] if i < len(images) else None)
        else:
            add_image_placeholder(0.5, y_start, 12.333, height, f"配图展示 ({image_count}张)")

    def _fill_image_slide(self, slide, config: dict, colors: dict) -> None:
        """填充图片页，支持多种布局。
        
        布局类型：
        - left_half: 左侧图片，右侧文字
        - right_half: 左侧文字，右侧图片
        - left_large: 左侧大图（约60%），右侧内容
        - right_large: 左侧内容，右侧大图（约60%）
        - full_bg: 全屏背景图，叠加文字
        - top_bottom: 上图下文
        """
        title = config.get("title", "")
        subtitle = config.get("subtitle", "")
        image_path = config.get("image_path", "")
        layout = config.get("image_layout", config.get("layout", "left_half"))
        content = config.get("content", [])

        # 顶部装饰
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors["primary"]
        top_bar.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        # 根据布局类型放置图片和内容
        if layout == "full_bg":
            # 全屏背景图
            if image_path and Path(image_path).exists():
                try:
                    # 添加背景图片
                    pic = slide.shapes.add_picture(
                        image_path,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5)
                    )
                    # 将图片移到最底层
                    spTree = slide.shapes._spTree
                    pic_elem = pic._element
                    spTree.remove(pic_elem)
                    spTree.insert(2, pic_elem)  # 插入到第2个位置（背景之后）
                except Exception as e:
                    logger.warning(f"插入背景图片失败: {e}")
            
            # 叠加半透明遮罩
            overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            overlay.fill.fore_color.brightness = 0.7  # 70%透明度
            overlay.line.fill.background()
            
        elif layout in ["left_half", "left_large"]:
            # 左侧图片
            img_width = 6.0 if layout == "left_half" else 7.5
            if image_path and Path(image_path).exists():
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(0.5), Inches(1.3),
                        Inches(img_width), Inches(5.5)
                    )
                except Exception as e:
                    logger.warning(f"插入左侧图片失败: {e}")
            
            # 右侧内容
            x_start = 0.5 + img_width + 0.3
            content_width = 13.333 - x_start - 0.5
            
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(x_start), Inches(1.5), Inches(content_width), Inches(0.5))
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(18)
                p.font.color.rgb = colors["accent"]
            
            y = 2.2 if subtitle else 1.8
            for i, item in enumerate(content):
                box = slide.shapes.add_textbox(Inches(x_start), Inches(y + i * 0.9), Inches(content_width), Inches(0.9))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = colors["secondary"]

        elif layout in ["right_half", "right_large"]:
            # 左侧内容
            img_width = 6.0 if layout == "right_half" else 7.5
            content_width = 13.333 - img_width - 1.5
            
            x_start = 0.5
            
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(x_start), Inches(1.5), Inches(content_width), Inches(0.5))
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(18)
                p.font.color.rgb = colors["accent"]
            
            y = 2.2 if subtitle else 1.8
            for i, item in enumerate(content):
                box = slide.shapes.add_textbox(Inches(x_start), Inches(y + i * 0.9), Inches(content_width), Inches(0.9))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = colors["secondary"]
            
            # 右侧图片
            img_x = x_start + content_width + 0.3
            if image_path and Path(image_path).exists():
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(img_x), Inches(1.3),
                        Inches(img_width - 0.3), Inches(5.5)
                    )
                except Exception as e:
                    logger.warning(f"插入右侧图片失败: {e}")

        elif layout == "top_bottom":
            # 上图下文
            if image_path and Path(image_path).exists():
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(0.5), Inches(1.3),
                        Inches(12.333), Inches(3.0)
                    )
                except Exception as e:
                    logger.warning(f"插入上方图片失败: {e}")
            
            # 下方文字
            for i, item in enumerate(content):
                box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5 + i * 0.7), Inches(12.333), Inches(0.7))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = colors["secondary"]

    def _generate_ppt(self, params: dict[str, Any]) -> ToolResult:
        """生成完整 PPT。"""
        topic = params.get("topic", "").strip()
        if not topic:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT 主题不能为空",
            )

        outline = params.get("outline", "").strip()
        style = params.get("style", "business")
        slide_count = int(params.get("slide_count", 8))
        output_filename = params.get("output_filename", "").strip()

        # 验证风格
        if style not in STYLE_COLORS:
            style = "business"

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9
            prs.slide_height = Inches(7.5)

            colors = self._get_style_colors(style)

            # 1. 标题页
            self._add_title_slide(prs, topic, colors)

            # 2. 内容页
            if outline:
                # 按大纲生成内容页
                points = [p.strip() for p in outline.split("\n") if p.strip()]
                for i, point in enumerate(points):
                    self._add_content_slide(
                        prs,
                        slide_title=point,
                        subtitle=f"第 {i + 1} 部分",
                        colors=colors,
                    )
            else:
                # 生成默认结构
                default_sections = [
                    "背景介绍",
                    "核心内容",
                    "关键要点",
                    "数据分析",
                    "案例展示",
                    "总结与展望",
                ]
                for i, section in enumerate(default_sections[: slide_count - 2]):
                    self._add_content_slide(
                        prs,
                        slide_title=section,
                        subtitle=f"第 {i + 1} 部分",
                        colors=colors,
                    )

            # 3. 感谢页
            self._add_thank_slide(prs, colors)

            # 生成文件名
            if not output_filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                # 清理主题中的特殊字符
                safe_topic = "".join(c for c in topic[:20] if c.isalnum() or c in " _-")
                output_filename = f"ppt_{safe_topic}_{timestamp}"

            output_path = self.output_dir / f"{output_filename}.pptx"
            prs.save(str(output_path))

            file_size = output_path.stat().st_size
            actual_slide_count = len(prs.slides)

            logger.info("PPT 生成成功: %s (%d 页)", output_path, actual_slide_count)

            return ToolResult(
                status=ToolResultStatus.SUCCESS,
                output=f"✅ PPT 生成完成\n📁 文件: {output_path.name}\n📊 大小: {file_size} 字节\n📄 共 {actual_slide_count} 页\n🎨 风格: {style}",
                data={
                    "file_path": str(output_path),
                    "file_name": output_path.name,
                    "file_size": file_size,
                    "slide_count": actual_slide_count,
                    "style": style,
                },
            )

        except Exception as e:
            logger.error("PPT 生成失败: %s", e, exc_info=True)
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT 生成失败: {e}",
            )

    def _add_slide(self, params: dict[str, Any]) -> ToolResult:
        """向已有 PPT 添加幻灯片。"""
        ppt_path = params.get("ppt_path", "").strip()
        title = params.get("title", "").strip()
        content = params.get("content", "").strip()
        layout = params.get("layout", "content")

        if not ppt_path:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT 文件路径不能为空",
            )

        if not title:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="幻灯片标题不能为空",
            )

        if not content:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="幻灯片内容不能为空",
            )

        ppt_file = Path(ppt_path)
        if not ppt_file.exists():
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT 文件不存在: {ppt_path}",
            )

        try:
            prs = Presentation(str(ppt_file))
            colors = self._get_style_colors("business")  # 默认使用商务风格

            # 解析内容为要点列表
            bullet_points = [p.strip() for p in content.split("\n") if p.strip()]

            if layout == "title":
                # 标题布局
                self._add_title_slide(prs, title, colors)
            elif layout == "blank":
                # 空白布局
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                self._set_slide_background(slide, colors["bg"])
                # 只添加标题
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
                )
                title_frame = title_box.text_frame
                p = title_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = colors["primary"]
            elif layout == "two_column":
                # 两列布局
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                self._set_slide_background(slide, colors["bg"])

                # 标题
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
                )
                title_frame = title_box.text_frame
                p = title_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = colors["primary"]

                # 分两列显示内容
                mid = len(bullet_points) // 2
                left_points = bullet_points[:mid] if mid > 0 else bullet_points
                right_points = bullet_points[mid:] if mid > 0 else []

                # 左列
                left_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.8), Inches(5.9), Inches(5.0)
                )
                left_frame = left_box.text_frame
                left_frame.word_wrap = True
                for i, point in enumerate(left_points):
                    if i == 0:
                        p = left_frame.paragraphs[0]
                    else:
                        p = left_frame.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = colors["secondary"]
                    p.space_after = Pt(10)

                # 右列
                right_box = slide.shapes.add_textbox(
                    Inches(6.9), Inches(1.8), Inches(5.9), Inches(5.0)
                )
                right_frame = right_box.text_frame
                right_frame.word_wrap = True
                for i, point in enumerate(right_points):
                    if i == 0:
                        p = right_frame.paragraphs[0]
                    else:
                        p = right_frame.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = colors["secondary"]
                    p.space_after = Pt(10)
            else:
                # 默认内容布局
                self._add_content_slide(
                    prs,
                    slide_title=title,
                    subtitle="",
                    colors=colors,
                    bullet_points=bullet_points,
                )

            # 保存回原文件
            prs.save(str(ppt_file))

            slide_count = len(prs.slides)
            logger.info("已向 PPT 添加幻灯片: %s (共 %d 页)", ppt_file, slide_count)

            return ToolResult(
                status=ToolResultStatus.SUCCESS,
                output=f"✅ 幻灯片添加成功\n📁 文件: {ppt_file.name}\n📄 当前共 {slide_count} 页\n📝 新增标题: {title}",
                data={
                    "file_path": str(ppt_file),
                    "file_name": ppt_file.name,
                    "slide_count": slide_count,
                    "new_slide_title": title,
                    "layout": layout,
                },
            )

        except Exception as e:
            logger.error("添加幻灯片失败: %s", e, exc_info=True)
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"添加幻灯片失败: {e}",
            )

    def _export_pdf(self, params: dict[str, Any]) -> ToolResult:
        """导出 PPT 为 PDF。"""
        ppt_path = params.get("ppt_path", "").strip()
        output_filename = params.get("output_filename", "").strip()

        if not ppt_path:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT 文件路径不能为空",
            )

        ppt_file = Path(ppt_path)
        if not ppt_file.exists():
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT 文件不存在: {ppt_path}",
            )

        # 生成输出文件名
        if not output_filename:
            output_filename = ppt_file.stem

        output_path = self.output_dir / f"{output_filename}.pdf"

        # 尝试使用 LibreOffice 转换
        libreoffice_path = shutil.which("libreoffice") or shutil.which("soffice")
        
        if libreoffice_path:
            try:
                result = subprocess.run(
                    [
                        libreoffice_path,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        str(self.output_dir),
                        str(ppt_file),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )

                if result.returncode == 0:
                    # LibreOffice 输出的文件名可能与输入同名
                    expected_output = self.output_dir / f"{ppt_file.stem}.pdf"
                    if expected_output.exists() and expected_output != output_path:
                        expected_output.rename(output_path)

                    if output_path.exists():
                        file_size = output_path.stat().st_size
                        logger.info("PDF 导出成功: %s", output_path)

                        return ToolResult(
                            status=ToolResultStatus.SUCCESS,
                            output=f"✅ PDF 导出成功\n📁 文件: {output_path.name}\n📊 大小: {file_size} 字节",
                            data={
                                "file_path": str(output_path),
                                "file_name": output_path.name,
                                "file_size": file_size,
                                "export_method": "libreoffice",
                            },
                        )
                else:
                    logger.warning("LibreOffice 转换失败: %s", result.stderr)
            except subprocess.TimeoutExpired:
                logger.warning("LibreOffice 转换超时")
            except Exception as e:
                logger.warning("LibreOffice 转换出错: %s", e)

        # LibreOffice 不可用或转换失败，返回手动导出提示
        return ToolResult(
            status=ToolResultStatus.SUCCESS,
            output=f"⚠️ 自动 PDF 导出需要安装 LibreOffice。\n请手动打开 PPT 后导出为 PDF。\n📁 PPT 文件: {ppt_file}",
            data={
                "ppt_path": str(ppt_file),
                "export_method": "manual",
                "suggestion": "请使用 PowerPoint 或 WPS 打开 PPT 文件后，选择「另存为」或「导出」为 PDF 格式。",
            },
        )

    def _add_chart_slide(self, params: dict[str, Any]) -> ToolResult:
        """添加图表幻灯片。"""
        ppt_path = params.get("ppt_path", "").strip()
        title = params.get("title", "").strip()
        chart_type = params.get("chart_type", "bar")
        chart_data_str = params.get("chart_data", "{}")
        style = params.get("style", "business")

        if not ppt_path or not title:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT路径和标题不能为空",
            )

        ppt_file = Path(ppt_path)
        if not ppt_file.exists():
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT文件不存在: {ppt_path}",
            )

        try:
            import json
            chart_data = json.loads(chart_data_str)

            prs = Presentation(str(ppt_file))
            colors = self._get_style_colors(style)

            # 创建图表
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE

            # 添加幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            self._set_slide_background(slide, colors["bg"])

            # 添加标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
            )
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = colors["primary"]

            # 创建图表数据
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get("labels", [])
            chart_data_obj.add_series("数据", chart_data.get("values", []))

            # 添加图表
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(10), Inches(5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED if chart_type == "bar"
                else XL_CHART_TYPE.PIE if chart_type == "pie"
                else XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_obj
            ).chart

            # 保存
            prs.save(str(ppt_file))
            logger.info("图表幻灯片添加成功: %s", ppt_file)

            return ToolResult(
                status=ToolResultStatus.SUCCESS,
                output=f"✅ 图表幻灯片添加成功\n📁 文件: {ppt_file.name}\n📊 图表类型: {chart_type}",
                data={
                    "file_path": str(ppt_file),
                    "chart_type": chart_type,
                    "title": title,
                },
            )

        except json.JSONDecodeError:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="图表数据格式错误，请使用JSON格式",
            )
        except Exception as e:
            logger.error("添加图表幻灯片失败: %s", e)
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"添加图表幻灯片失败: {e}",
            )

    def _add_table_slide(self, params: dict[str, Any]) -> ToolResult:
        """添加表格幻灯片。"""
        ppt_path = params.get("ppt_path", "").strip()
        title = params.get("title", "").strip()
        table_data_str = params.get("table_data", "[]")
        style = params.get("style", "business")

        if not ppt_path or not title:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT路径和标题不能为空",
            )

        ppt_file = Path(ppt_path)
        if not ppt_file.exists():
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT文件不存在: {ppt_path}",
            )

        try:
            import json
            table_data = json.loads(table_data_str)

            prs = Presentation(str(ppt_file))
            colors = self._get_style_colors(style)

            # 添加幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            self._set_slide_background(slide, colors["bg"])

            # 添加标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.0)
            )
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = colors["primary"]

            # 解析表格数据
            if isinstance(table_data, list) and len(table_data) > 0:
                if isinstance(table_data[0], dict) and "headers" in table_data[0]:
                    # 新格式
                    headers = table_data[0].get("headers", [])
                    rows = [item.get("row", []) for item in table_data[1:] if isinstance(item, dict)]
                    all_rows = [headers] + rows
                else:
                    all_rows = table_data

                if all_rows and len(all_rows) > 0:
                    num_rows = len(all_rows)
                    num_cols = max(len(row) for row in all_rows) if all_rows else 0

                    # 添加表格
                    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12), Inches(5)
                    table = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy).table

                    # 填充数据
                    for i, row in enumerate(all_rows):
                        for j, cell_value in enumerate(row):
                            if j < num_cols:
                                cell = table.cell(i, j)
                                cell.text = str(cell_value) if cell_value else ""
                                # 表头样式
                                if i == 0:
                                    for para in cell.text_frame.paragraphs:
                                        para.font.bold = True
                                        para.font.size = Pt(14)
                                        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = colors["primary"]
                                else:
                                    for para in cell.text_frame.paragraphs:
                                        para.font.size = Pt(12)
                                        para.font.color.rgb = colors["secondary"]

            # 保存
            prs.save(str(ppt_file))
            logger.info("表格幻灯片添加成功: %s", ppt_file)

            return ToolResult(
                status=ToolResultStatus.SUCCESS,
                output=f"✅ 表格幻灯片添加成功\n📁 文件: {ppt_file.name}\n📊 表格数据已添加",
                data={
                    "file_path": str(ppt_file),
                    "title": title,
                },
            )

        except json.JSONDecodeError:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="表格数据格式错误，请使用JSON格式",
            )
        except Exception as e:
            logger.error("添加表格幻灯片失败: %s", e)
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"添加表格幻灯片失败: {e}",
            )

    def _add_section_divider(self, params: dict[str, Any]) -> ToolResult:
        """添加章节分隔页。"""
        ppt_path = params.get("ppt_path", "").strip()
        section_title = params.get("section_title", "").strip()
        section_number = params.get("section_number", 1)
        style = params.get("style", "business")

        if not ppt_path or not section_title:
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error="PPT路径和章节标题不能为空",
            )

        ppt_file = Path(ppt_path)
        if not ppt_file.exists():
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"PPT文件不存在: {ppt_path}",
            )

        try:
            prs = Presentation(str(ppt_file))
            colors = self._get_style_colors(style)

            # 添加幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            self._set_slide_background(slide, colors["primary"])

            # 章节编号
            num_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
            )
            num_frame = num_box.text_frame
            p = num_frame.paragraphs[0]
            p.text = f"第 {section_number} 部分"
            p.font.size = Pt(24)
            p.font.color.rgb = colors["accent"]
            p.alignment = PP_ALIGN.CENTER

            # 章节标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.5), Inches(12.333), Inches(2)
            )
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            # 保存
            prs.save(str(ppt_file))
            logger.info("章节分隔页添加成功: %s", ppt_file)

            return ToolResult(
                status=ToolResultStatus.SUCCESS,
                output=f"✅ 章节分隔页添加成功\n📁 文件: {ppt_file.name}\n📋 第 {section_number} 部分: {section_title}",
                data={
                    "file_path": str(ppt_file),
                    "section_number": section_number,
                    "section_title": section_title,
                },
            )

        except Exception as e:
            logger.error("添加章节分隔页失败: %s", e)
            return ToolResult(
                status=ToolResultStatus.ERROR,
                error=f"添加章节分隔页失败: {e}",
            )


# 用于测试
if __name__ == "__main__":
    import asyncio

    async def test():
        tool = PPTTool()
        
        # 测试生成 PPT
        result = await tool.execute(
            "generate_ppt",
            {
                "topic": "人工智能发展趋势",
                "outline": "AI 历史回顾\n当前技术突破\n应用场景分析\n未来发展预测",
                "style": "business",
            },
        )
        print("生成 PPT 结果:")
        print(result.output)
        print("Data:", result.data)
        
        if result.is_success:
            # 测试添加幻灯片
            ppt_path = result.data["file_path"]
            result2 = await tool.execute(
                "add_slide",
                {
                    "ppt_path": ppt_path,
                    "title": "补充说明",
                    "content": "第一个要点\n第二个要点\n第三个要点",
                    "layout": "content",
                },
            )
            print("\n添加幻灯片结果:")
            print(result2.output)

    asyncio.run(test())

    # ==================== 增强布局方法 ====================

    def add_decorative_line(self, slide, x: float, y: float, width: float, color: RGBColor) -> None:
        """添加装饰线条"""
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        
        line_xml = parse_xml(
            f'<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:solidFill><a:srgbClr val="{color[0]:02X}{color[1]:02X}{color[2]:02X}"/></a:solidFill>'
            f'</a:ln>'
        )
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(0.03))
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = color
        shape.fill.background()

    def add_icon_card(self, slide, icon: str, title: str, description: str, 
                      x: float, y: float, width: float, height: float, 
                      colors: dict) -> None:
        """添加图标卡片"""
        # 卡片背景
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors.get("bg", RGBColor(0xFF, 0xFF, 0xFF))
        shape.line.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        shape.line.width = Pt(1)

        # 图标
        icon_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(0.8), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = ICONS_MAP.get(icon, icon)
        p.font.size = Pt(32)
        p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

        # 标题
        title_box = slide.shapes.add_textbox(Inches(x + 1.0), Inches(y + 0.2), Inches(width - 1.2), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(width - 0.4), Inches(height - 0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(11)
        p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))

    def add_image_grid(self, slide, image_count: int, colors: dict, 
                       y_start: float = 1.8, height: float = 4.5) -> None:
        """添加图片网格占位区域（带装饰性边框）"""
        if image_count == 1:
            # 单图大布局
            self._add_image_placeholder(slide, 0.5, y_start, 12.333, height, colors, "主图")
        elif image_count == 2:
            # 左右两图
            self._add_image_placeholder(slide, 0.5, y_start, 5.8, height, colors, "图片 1")
            self._add_image_placeholder(slide, 6.8, y_start, 5.8, height, colors, "图片 2")
        elif image_count == 3:
            # 左大右小（上下两图）
            self._add_image_placeholder(slide, 0.5, y_start, 6.0, height, colors, "图片 1")
            self._add_image_placeholder(slide, 6.8, y_start, 5.8, height / 2 - 0.1, colors, "图片 2")
            self._add_image_placeholder(slide, 6.8, y_start + height / 2 + 0.1, 5.8, height / 2 - 0.1, colors, "图片 3")
        elif image_count == 4:
            # 2x2 网格
            w = 5.8
            h = height / 2 - 0.15
            gap = 0.2
            positions = [
                (0.5, y_start),
                (0.5 + w + gap, y_start),
                (0.5, y_start + h + gap),
                (0.5 + w + gap, y_start + h + gap)
            ]
            for i, (x, y) in enumerate(positions):
                self._add_image_placeholder(slide, x, y, w, h, colors, f"图片 {i+1}")
        elif image_count == 5:
            # 上3下2
            w1 = 3.8
            w2 = 5.8
            h_small = height / 2 - 0.15
            positions = [
                (0.5, y_start, w1, h_small),
                (0.5 + w1 + 0.2, y_start, w1, h_small),
                (0.5 + (w1 + 0.2) * 2, y_start, w1, h_small),
                (0.5, y_start + h_small + 0.3, w2, h_small),
                (0.5 + w2 + 0.4, y_start + h_small + 0.3, w2, h_small)
            ]
            for i, (x, y, w, h) in enumerate(positions):
                self._add_image_placeholder(slide, x, y, w, h, colors, f"图片 {i+1}")

    def _add_image_placeholder(self, slide, x: float, y: float, width: float, height: float, 
                                colors: dict, label: str) -> None:
        """添加图片占位区域"""
        # 背景框
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors.get("highlight", RGBColor(0xE3, 0xF2, 0xFD))
        shape.line.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        shape.line.width = Pt(1.5)

        # 中心图标
        icon_box = slide.shapes.add_textbox(Inches(x), Inches(y + height/2 - 0.5), Inches(width), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = "🖼️"
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

        # 标签
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y + height/2 + 0.3), Inches(width), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))

    def add_icon_header(self, slide, title: str, subtitle: str, colors: dict) -> None:
        """添加带图标装饰的标题头"""
        # 顶部装饰线
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.3), Inches(0.8), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        line.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

        # 底部装饰线
        line2 = slide.shapes.add_shape(1, Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.02))
        line2.fill.solid()
        line2.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
        line2.line.fill.background()

    def add_icon_list_content(self, slide, items: list, icons: list, colors: dict,
                                y_start: float = 2.0, col: int = 1) -> None:
        """添加带图标的列表内容"""
        if col == 1:
            # 单栏
            for i, item in enumerate(items):
                y = y_start + i * 0.9
                # 图标
                icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.6), Inches(0.6))
                tf = icon_box.text_frame
                p = tf.paragraphs[0]
                icon = icons[i] if i < len(icons) else "●"
                p.text = ICONS_MAP.get(icon, icon)
                p.font.size = Pt(24)
                p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

                # 内容
                content_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.0), Inches(0.8))
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))
        else:
            # 双栏
            mid = len(items) // 2
            left_items = items[:mid]
            right_items = items[mid:]

            for i, item in enumerate(left_items):
                y = y_start + i * 0.9
                icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.5), Inches(0.5))
                tf = icon_box.text_frame
                p = tf.paragraphs[0]
                icon = icons[i] if i < len(icons) else "●"
                p.text = ICONS_MAP.get(icon, icon)
                p.font.size = Pt(20)
                p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

                content_box = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(5.5), Inches(0.8))
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))

            for i, item in enumerate(right_items):
                y = y_start + i * 0.9
                idx = i + mid
                icon_box = slide.shapes.add_textbox(Inches(6.8), Inches(y), Inches(0.5), Inches(0.5))
                tf = icon_box.text_frame
                p = tf.paragraphs[0]
                icon = icons[idx] if idx < len(icons) else "●"
                p.text = ICONS_MAP.get(icon, icon)
                p.font.size = Pt(20)
                p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

                content_box = slide.shapes.add_textbox(Inches(7.4), Inches(y), Inches(5.5), Inches(0.8))
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))

    def add_stunning_title_slide(self, prs: Presentation, topic: str, subtitle: str, 
                                  colors: dict) -> None:
        """生成惊艳的标题页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors.get("bg", RGBColor(0xF0, 0xF8, 0xFF)))

        # 左侧装饰条
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        bar.line.fill.background()

        # 右上角装饰圆
        circle = slide.shapes.add_shape(9, Inches(11), Inches(-1), Inches(3), Inches(3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
        circle.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(10), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

        # 标题下划线
        line = slide.shapes.add_shape(1, Inches(0.8), Inches(3.7), Inches(3), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))
        line.line.fill.background()

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(10), Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))

        # 日期
        date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(5), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(14)
        p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))
