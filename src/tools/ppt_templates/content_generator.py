"""PPT 内容生成器 - 生成图表、表格等丰富内容。

支持：
- 柱状图 (bar)
- 饼图 (pie)
- 折线图 (line)
- 表格
"""

import logging
from typing import Any
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)


class RichContentGenerator:
    """PPT 丰富内容生成器"""

    # 图表类型映射
    CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
        "bar_3d": XL_CHART_TYPE.BAR_CLUSTERED,  # 降级为普通柱状图
        "pie": XL_CHART_TYPE.PIE,
        "pie_3d": XL_CHART_TYPE.PIE,  # 降级为普通饼图
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "line_3d": XL_CHART_TYPE.LINE,  # 降级为普通折线图
        "area": XL_CHART_TYPE.AREA,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }

    def __init__(self, image_tool=None):
        """初始化内容生成器

        Args:
            image_tool: 图片生成工具实例（可选）
        """
        self.image_tool = image_tool

    def create_chart_data(self, chart_type: str, data: dict[str, Any]) -> CategoryChartData:
        """创建图表数据对象

        Args:
            chart_type: 图表类型 (bar/pie/line)
            data: 数据字典，包含 labels 和 values

        Returns:
            CategoryChartData 对象
        """
        chart_data = CategoryChartData()
        chart_data.categories = data.get("labels", [])
        chart_data.add_series(
            data.get("series_name", "数据"),
            data.get("values", [])
        )
        return chart_data

    def get_chart_type(self, chart_type: str) -> XL_CHART_TYPE:
        """获取图表类型枚举

        Args:
            chart_type: 图表类型字符串

        Returns:
            XL_CHART_TYPE 枚举值
        """
        chart_type_lower = chart_type.lower().strip()
        return self.CHART_TYPE_MAP.get(
            chart_type_lower,
            XL_CHART_TYPE.BAR_CLUSTERED  # 默认柱状图
        )

    def add_chart_to_slide(
        self,
        slide,
        chart_type: str,
        data: dict[str, Any],
        title: str = "",
        left: float = 1.0,
        top: float = 1.8,
        width: float = 11.0,
        height: float = 4.5,
        style_colors: dict = None,
    ) -> Any:
        """向幻灯片添加图表

        Args:
            slide: 幻灯片对象
            chart_type: 图表类型
            data: 图表数据
            title: 图表标题
            left, top, width, height: 图表位置和大小
            style_colors: 样式颜色配置

        Returns:
            图表对象
        """
        try:
            chart_data = self.create_chart_data(chart_type, data)
            xl_chart_type = self.get_chart_type(chart_type)

            chart = slide.shapes.add_chart(
                xl_chart_type,
                Inches(left), Inches(top), Inches(width), Inches(height),
                chart_data
            ).chart

            # 设置图表标题
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.paragraphs[0].text = title
                if style_colors:
                    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = style_colors.get("primary")

            # 设置图例
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            # 应用配色
            if style_colors:
                try:
                    from pptx.dml.color import RGBColor
                    plot = chart.plots[0]
                    if hasattr(plot, "series"):
                        for series in plot.series:
                            series.format.fill.solid()
                            series.format.fill.fore_color.rgb = style_colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
                except Exception as e:
                    logger.warning(f"图表配色设置失败: {e}")

            return chart

        except Exception as e:
            logger.error(f"添加图表失败: {e}")
            raise

    def add_table_to_slide(
        self,
        slide,
        headers: list[str],
        rows: list[list[str]],
        title: str = "",
        left: float = 0.5,
        top: float = 1.8,
        width: float = 12.0,
        height: float = 4.5,
        style_colors: dict = None,
    ) -> Any:
        """向幻灯片添加表格

        Args:
            slide: 幻灯片对象
            headers: 表头列表
            rows: 数据行列表
            title: 表格标题
            left, top, width, height: 表格位置和大小
            style_colors: 样式颜色配置

        Returns:
            表格对象
        """
        try:
            num_rows = len(rows) + 1  # +1 表头
            num_cols = max(len(headers), max(len(row) for row in rows) if rows else 0)

            if num_cols == 0:
                return None

            # 添加表格
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            table = table_shape.table

            # 设置列宽
            col_width = Emu(int(Inches(width).emu / num_cols))
            for col_idx in range(num_cols):
                table.columns[col_idx].width = col_width

            # 填充表头
            for col_idx, header in enumerate(headers[:num_cols]):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                # 表头样式
                self._style_table_cell(
                    cell,
                    is_header=True,
                    style_colors=style_colors
                )

            # 填充数据行
            for row_idx, row in enumerate(rows):
                for col_idx, value in enumerate(row[:num_cols]):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(value)
                    self._style_table_cell(
                        cell,
                        is_header=False,
                        style_colors=style_colors
                    )

            return table

        except Exception as e:
            logger.error(f"添加表格失败: {e}")
            raise

    def _style_table_cell(self, cell, is_header: bool = False, style_colors: dict = None) -> None:
        """设置表格单元格样式

        Args:
            cell: 单元格对象
            is_header: 是否表头
            style_colors: 样式颜色配置
        """
        from pptx.dml.color import RGBColor

        # 设置文字
        for para in cell.text_frame.paragraphs:
            para.alignment = 1  # 居中对齐
            for run in para.runs:
                run.font.size = Pt(12 if is_header else 11)
                run.font.bold = is_header
                if style_colors:
                    run.font.color.rgb = style_colors.get("text", RGBColor(0x33, 0x33, 0x33))

        # 设置填充色
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = style_colors.get("primary", RGBColor(0x00, 0x52, 0x8A)) if style_colors else RGBColor(0x00, 0x52, 0x8A)
            # 表头文字设为白色
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            cell.fill.fore_color.rgb = style_colors.get("bg", RGBColor(0xFF, 0xFF, 0xFF)) if style_colors else RGBColor(0xFF, 0xFF, 0xFF)

        # 设置边框
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml

        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            # 简单的边框设置
            lnL = parse_xml(
                r'<a:lnL w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                r'<a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:lnL>'
            )
            lnR = parse_xml(
                r'<a:lnR w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                r'<a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:lnR>'
            )
            lnT = parse_xml(
                r'<a:lnT w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                r'<a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:lnT>'
            )
            lnB = parse_xml(
                r'<a:lnB w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                r'<a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill></a:lnB>'
            )
            tcPr.append(lnL)
            tcPr.append(lnR)
            tcPr.append(lnT)
            tcPr.append(lnB)
        except Exception as e:
            logger.debug(f"表格边框设置失败: {e}")

    async def generate_image(self, prompt: str, output_dir: Path = None) -> str | None:
        """生成配图

        Args:
            prompt: 图片生成提示词
            output_dir: 输出目录

        Returns:
            生成的图片路径，失败返回 None
        """
        if not self.image_tool:
            logger.debug("未配置图片生成工具，跳过配图生成")
            return None

        try:
            # 调用图片生成工具
            result = await self.image_tool.execute(
                "generate_image",
                {
                    "prompt": prompt,
                    "output_dir": str(output_dir) if output_dir else "",
                    "size": "16:9",
                }
            )

            if result.is_success and result.data:
                return result.data.get("file_path")

            logger.warning(f"图片生成失败: {result.error if hasattr(result, 'error') else '未知错误'}")
            return None

        except Exception as e:
            logger.error(f"图片生成异常: {e}")
            return None

    def create_two_column_content(
        self,
        slide,
        left_title: str,
        left_content: list[str],
        right_title: str,
        right_content: list[str],
        style_colors: dict = None,
    ) -> None:
        """创建双栏内容

        Args:
            slide: 幻灯片对象
            left_title: 左侧标题
            left_content: 左侧内容列表
            right_title: 右侧标题
            right_content: 右侧内容列表
            style_colors: 样式颜色配置
        """
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # 左侧内容
        self._add_content_box(
            slide,
            left_title,
            left_content,
            left=0.5,
            top=1.8,
            width=5.9,
            height=4.5,
            style_colors=style_colors,
        )

        # 右侧内容
        self._add_content_box(
            slide,
            right_title,
            right_content,
            left=6.9,
            top=1.8,
            width=5.9,
            height=4.5,
            style_colors=style_colors,
        )

    def _add_content_box(
        self,
        slide,
        title: str,
        content: list[str],
        left: float,
        top: float,
        width: float,
        height: float,
        style_colors: dict = None,
    ) -> None:
        """添加内容文本框"""
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        if style_colors:
            p.font.color.rgb = style_colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + 0.6), Inches(width), Inches(height - 0.6)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content):
            p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            if style_colors:
                p.font.color.rgb = style_colors.get("text", RGBColor(0x33, 0x33, 0x33))
            p.space_after = Pt(8)
