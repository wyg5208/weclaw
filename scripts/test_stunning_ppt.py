# -*- coding: utf-8 -*-
"""生成惊艳的PPT测试脚本"""
import asyncio
import sys
sys.path.insert(0, '.')

from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.tools.ppt_generator import PPTTool, STYLE_COLORS_ENHANCED, ICONS_MAP


def get_gradient_colors():
    """获取渐变配色方案"""
    schemes = list(STYLE_COLORS_ENHANCED.values())
    return schemes


def add_stunning_title_slide(prs, topic, subtitle, colors):
    """生成惊艳标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = colors.get("bg", RGBColor(0xF0, 0xF8, 0xFF))

    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
    bar.line.fill.background()

    # 右上角装饰圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    circle.line.fill.background()

    # 左下角装饰
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(6), Inches(2), Inches(2))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))
    circle2.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    circle2.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

    # 标题下划线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6), Inches(4), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))
    line.line.fill.background()

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(10), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))

    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(5), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%Y年%m月%d日")
    p.font.size = Pt(14)
    p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))


def add_content_slide_with_icons(prs, title, subtitle, items, icons, colors):
    """添加带图标装饰的内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = colors.get("bg", RGBColor(0xFF, 0xFF, 0xFF))

    # 顶部装饰
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
    top_bar.line.fill.background()

    # 标题左侧图标
    icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = "★"
    p.font.size = Pt(28)
    p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.0), Inches(11), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    line.line.fill.background()

    # 图标列表内容
    y_start = 1.8
    for i, item in enumerate(items):
        y = y_start + i * 1.0
        
        # 圆点背景
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y + 0.1), Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        dot.line.fill.background()

        # 图标
        icon_char = ICONS_MAP.get(icons[i] if i < len(icons) else "check", "✓")
        icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.05), Inches(0.35), Inches(0.4))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon_char
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11.5), Inches(0.9))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))


def add_image_grid_slide(prs, title, image_count, colors):
    """添加图片网格布局页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = colors.get("bg", RGBColor(0xFF, 0xFF, 0xFF))

    # 标题
    icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🖼️"
    p.font.size = Pt(28)

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

    # 图片网格
    y_start = 1.3
    height = 5.5
    
    if image_count == 1:
        _add_image_placeholder_styled(slide, 0.5, y_start, 12.333, height, colors, "核心配图")
    elif image_count == 2:
        _add_image_placeholder_styled(slide, 0.5, y_start, 6.0, height, colors, "配图 1")
        _add_image_placeholder_styled(slide, 6.8, y_start, 6.0, height, colors, "配图 2")
    elif image_count == 3:
        _add_image_placeholder_styled(slide, 0.5, y_start, 6.0, height, colors, "配图 1")
        _add_image_placeholder_styled(slide, 6.8, y_start, 5.8, height/2 - 0.1, colors, "配图 2")
        _add_image_placeholder_styled(slide, 6.8, y_start + height/2 + 0.1, 5.8, height/2 - 0.1, colors, "配图 3")
    elif image_count == 4:
        w = 6.0
        h = height/2 - 0.1
        gap = 0.2
        positions = [(0.5, y_start), (0.5 + w + gap, y_start), (0.5, y_start + h + gap), (0.5 + w + gap, y_start + h + gap)]
        for i, (x, y) in enumerate(positions):
            _add_image_placeholder_styled(slide, x, y, w, h, colors, f"配图 {i+1}")
    elif image_count == 6:
        # 3x2 网格
        w = 3.9
        h = 2.5
        gap = 0.2
        positions = [
            (0.5, y_start), (0.5 + w + gap, y_start), (0.5 + (w + gap) * 2, y_start),
            (0.5, y_start + h + gap), (0.5 + w + gap, y_start + h + gap), (0.5 + (w + gap) * 2, y_start + h + gap)
        ]
        for i, (x, y) in enumerate(positions):
            _add_image_placeholder_styled(slide, x, y, w, h, colors, f"配图 {i+1}")


def _add_image_placeholder_styled(slide, x, y, w, h, colors, label):
    """带装饰的图片占位区域"""
    # 阴影边框效果
    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(y + 0.05), Inches(w), Inches(h))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = colors.get("secondary", RGBColor(0xCC, 0xCC, 0xCC))
    shadow.line.fill.background()

    # 主框
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors.get("highlight", RGBColor(0xE3, 0xF2, 0xFD))
    shape.line.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
    shape.line.width = Pt(2)

    # 中心图标
    icon_box = slide.shapes.add_textbox(Inches(x), Inches(y + h/2 - 0.6), Inches(w), Inches(1))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📷"
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

    # 标签
    label_box = slide.shapes.add_textbox(Inches(x), Inches(y + h/2 + 0.4), Inches(w), Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))


def add_icon_cards_slide(prs, title, cards_data, colors):
    """添加图标卡片布局页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = colors.get("bg", RGBColor(0xFF, 0xFF, 0xFF))

    # 标题
    icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯"
    p.font.size = Pt(28)
    p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    line.line.fill.background()

    # 卡片数量和布局
    card_count = len(cards_data)
    if card_count == 3:
        w, h = 3.9, 5.5
        gap = 0.3
        positions = [(0.5, 1.4), (0.5 + w + gap, 1.4), (0.5 + (w + gap) * 2, 1.4)]
    elif card_count == 4:
        w, h = 5.8, 2.5
        gap = 0.3
        positions = [(0.5, 1.4), (0.5 + w + gap, 1.4), (0.5, 1.4 + h + gap), (0.5 + w + gap, 1.4 + h + gap)]
    elif card_count == 6:
        w, h = 3.9, 2.5
        gap = 0.2
        positions = [
            (0.5, 1.4), (0.5 + w + gap, 1.4), (0.5 + (w + gap) * 2, 1.4),
            (0.5, 1.4 + h + gap), (0.5 + w + gap, 1.4 + h + gap), (0.5 + (w + gap) * 2, 1.4 + h + gap)
        ]
    else:
        w, h = 5.8, 5.5
        positions = [(0.5, 1.4)]

    for i, card in enumerate(cards_data):
        if i >= len(positions):
            break
        x, y = positions[i]
        icon = card.get("icon", "★")
        card_title = card.get("title", f"要点 {i+1}")
        card_desc = card.get("description", "")

        # 卡片背景
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        shape.line.width = Pt(1.5)

        # 顶部装饰条
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.15))
        top.fill.solid()
        top.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
        top.line.fill.background()

        # 图标
        icon_char = ICONS_MAP.get(icon, icon)
        icon_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.4), Inches(w), Inches(1))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon_char
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

        # 标题
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.5), Inches(w - 0.4), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

        # 描述
        if card_desc:
            desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.1), Inches(w - 0.4), Inches(h - 2.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card_desc
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))


def add_chart_slide_styled(prs, title, chart_data, colors):
    """添加带装饰的图表页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = colors.get("bg", RGBColor(0xFF, 0xFF, 0xFF))

    # 标题
    icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(0.6), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊"
    p.font.size = Pt(28)
    p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.3), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    line.line.fill.background()

    # 图表
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data.get("labels", ["A", "B", "C", "D"])
    chart_data_obj.add_series("数据", chart_data.get("values", [75, 85, 60, 90]))

    x, y, cx, cy = Inches(1), Inches(1.5), Inches(10), Inches(5.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data_obj).chart
    chart.has_legend = False
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))


def add_stunning_thank_slide(prs, colors):
    """生成惊艳的感谢页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = colors.get("bg", RGBColor(0xF0, 0xF8, 0xFF))

    # 左侧装饰
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))
    bar.line.fill.background()

    # 右上装饰
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(3.5), Inches(3.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    circle.line.fill.background()

    # 主文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(72)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = colors.get("primary", RGBColor(0x00, 0x52, 0x8A))

    # 英文
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = colors.get("accent", RGBColor(0x00, 0x96, 0xD6))

    # 下划线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.0), Inches(4.333), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = colors.get("highlight", RGBColor(0x4F, 0xC3, 0xF7))
    line.line.fill.background()

    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "期待与您深入交流"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = colors.get("secondary", RGBColor(0x33, 0x33, 0x33))


async def generate_stunning_ppt():
    """生成惊艳的PPT"""
    print("\n" + "="*60)
    print("  生成惊艳PPT测试")
    print("="*60 + "\n")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 使用蓝色渐变配色
    colors = STYLE_COLORS_ENHANCED["gradient_blue"]

    # 第1页：惊艳标题页
    print("  [1/10] 生成惊艳标题页...")
    add_stunning_title_slide(prs, "人工智能发展趋势报告", "引领未来的技术创新与产业变革", colors)

    # 第2页：目录页
    print("  [2/10] 生成目录页...")
    add_content_slide_with_icons(
        prs, "内容概览", "本次报告核心章节",
        [
            "AI 技术演进：从机器学习到深度学习的突破",
            "技术突破：大模型与生成式 AI 的崛起",
            "应用场景：各行业的 AI 落地实践",
            "未来展望：AGI 与人类协作的新时代",
        ],
        ["chart", "rocket", "globe", "light"], colors
    )

    # 第3页：图片网格 - 2图
    print("  [3/10] 生成双图布局页...")
    add_image_grid_slide(prs, "AI 技术发展历程", 2, colors)

    # 第4页：图标卡片 - 4卡片
    print("  [4/10] 生成四卡片布局页...")
    add_icon_cards_slide(prs, "核心能力矩阵", [
        {"icon": "tech", "title": "深度学习", "description": "神经网络模型持续优化，参数规模突破万亿级别"},
        {"icon": "chart", "title": "数据处理", "description": "海量数据清洗与标注，支撑模型训练需求"},
        {"icon": "globe", "title": "跨模态理解", "description": "文本、图像、视频的统一表示学习"},
        {"icon": "rocket", "title": "实时推理", "description": "边缘部署与模型压缩技术突破"},
    ], colors)

    # 第5页：图表页
    print("  [5/10] 生成图表页...")
    add_chart_slide_styled(prs, "AI 市场增长趋势", {
        "labels": ["2022", "2023", "2024", "2025", "2026"],
        "values": [450, 620, 890, 1250, 1680]
    }, colors)

    # 第6页：图片网格 - 4图
    print("  [6/10] 生成四图网格页...")
    add_image_grid_slide(prs, "行业应用场景", 4, colors)

    # 第7页：图标卡片 - 3卡片
    print("  [7/10] 生成三卡片布局页...")
    add_icon_cards_slide(prs, "技术创新方向", [
        {"icon": "star", "title": "Transformer 架构", "description": "自注意力机制成为主流，Attention is All You Need"},
        {"icon": "fire", "title": "大模型时代", "description": "GPT-4、Claude、Gemini 等百亿参数模型涌现"},
        {"icon": "target", "title": "Agent 智能体", "description": "自主规划、工具调用、多智能体协作"},
    ], colors)

    # 第8页：内容页 - 双栏
    print("  [8/10] 生成双栏内容页...")
    add_content_slide_with_icons(
        prs, "典型应用案例", "AI 赋能千行百业",
        [
            "智能客服：自然语言理解与对话生成",
            "内容创作：文章、代码、图像、视频生成",
            "数据分析：智能洞察与决策支持",
            "代码开发：Copilot 提升编程效率",
            "医疗诊断：影像识别与辅助诊疗",
            "自动驾驶：环境感知与路径规划",
        ],
        ["people", "light", "chart", "tech", "target", "rocket"], colors
    )

    # 第9页：图片网格 - 6图
    print("  [9/10] 生成六图网格页...")
    add_image_grid_slide(prs, "生态合作伙伴", 6, colors)

    # 第10页：感谢页
    print("  [10/10] 生成感谢页...")
    add_stunning_thank_slide(prs, colors)

    # 保存
    output_dir = Path("generated")
    output_dir.mkdir(exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = output_dir / f"stunning_ppt_惊艳版_{timestamp}.pptx"
    prs.save(str(output_path))

    file_size = output_path.stat().st_size
    slide_count = len(prs.slides)

    print("\n" + "="*60)
    print(f"  PPT 生成成功！")
    print("="*60)
    print(f"  文件: {output_path.name}")
    print(f"  大小: {file_size:,} 字节")
    print(f"  页数: {slide_count} 页")
    print(f"  配色: 渐变蓝")
    print("="*60)
    
    return str(output_path)


if __name__ == "__main__":
    result = asyncio.run(generate_stunning_ppt())
    print(f"\n请打开查看: {result}")
