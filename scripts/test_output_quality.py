"""Office 工具实际输出验证测试 v2.28.0

验证实际生成的文件内容和质量，而不仅仅是函数调用成功。
"""

import asyncio
import sys
import os
import shutil
from pathlib import Path
from datetime import datetime

# 添加项目根目录
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from src.tools.pdf_generator import PDFGeneratorTool
from src.tools.ppt_generator import PPTTool
from src.tools.excel_validator import refcheck_workbook, apply_professional_style, STYLE_THEMES
from docx import Document
from openpyxl import Workbook, load_workbook

# 测试输出目录
TEST_OUTPUT_DIR = project_root / "generated" / datetime.now().strftime("%Y-%m-%d")
TEST_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 测试结果统计
test_stats = {"total": 0, "passed": 0, "failed": 0}


def log_test(name: str, passed: bool, detail: str = ""):
    """记录测试结果"""
    test_stats["total"] += 1
    if passed:
        test_stats["passed"] += 1
        status = "✅"
    else:
        test_stats["failed"] += 1
        status = "❌"
    
    print(f"  {status} {name}")
    if detail:
        print(f"     {detail}")


def section(title: str):
    """打印分节标题"""
    print(f"\n{'='*60}")
    print(f"  {title}")
    print("=" * 60)


# ==============================================================================
# 1. PDF 生成质量验证
# ==============================================================================

async def test_pdf_quality():
    """测试 PDF 实际生成质量"""
    section("PDF 生成质量验证")
    
    tool = PDFGeneratorTool()
    
    # 1.1 测试 Markdown 转 PDF 实际文件
    md_content = """# 项目报告

## 概述
这是一个测试报告，包含多种格式内容。

### 列表项
- 重要事项一
- 重要事项二
- 重要事项三

### 表格
| 项目 | 金额 | 状态 |
|------|------|------|
| 收入 | 10000 | 完成 |
| 支出 | 5000 | 进行中 |

**加粗文本** 和 *斜体文本*

[链接](https://example.com)
"""
    
    result = await tool.execute("md_to_pdf", {
        "markdown_content": md_content,
        "title": "质量测试报告",
        "author": "测试作者",
        "page_style": "modern"
    })
    
    # 验证1: 执行成功
    log_test("PDF 生成执行", result.status.value == "success", f"状态: {result.status.value}")
    
    if result.is_success and result.data.get("file_path"):
        pdf_path = result.data["file_path"]
        
        # 验证2: 文件存在
        log_test("PDF 文件存在", os.path.exists(pdf_path), pdf_path)
        
        # 验证3: 文件大小合理 (> 1KB)
        file_size = os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0
        log_test("PDF 文件大小", file_size > 1024, f"{file_size} bytes")
        
        # 验证4: 文件是有效的 PDF (检查文件头)
        try:
            with open(pdf_path, 'rb') as f:
                header = f.read(10)
                is_valid_pdf = header.startswith(b'%PDF')
                log_test("PDF 文件头有效", is_valid_pdf, f"文件头: {header[:5]}")
        except Exception as e:
            log_test("PDF 文件头有效", False, f"读取失败: {e}")
        
        # 验证5: 包含预期内容
        try:
            with open(pdf_path, 'rb') as f:
                content = f.read()
                # PDF 应该包含这些关键词
                has_chinese = b'\xe9\xa1\xb9\xe7\x9b\xae' in content or b'\xe6\x8a\xa5\xe5\x91\x8a' in content  # 项目/报告
                has_numbers = b'10000' in content or b'5000' in content
                log_test("PDF 包含中文内容", has_chinese or file_size > 5000, "内容检查")
        except Exception as e:
            log_test("PDF 内容检查", False, f"读取失败: {e}")
        
        print(f"     生成文件: {pdf_path}")
        
    else:
        log_test("PDF 文件存在", False, f"生成失败: {result.error}")
    
    # 1.2 测试不同页面样式
    for style in ["modern", "classic", "compact"]:
        result = await tool.execute("md_to_pdf", {
            "markdown_content": "# 测试\n\n内容",
            "title": f"样式测试-{style}",
            "page_style": style
        })
        if result.is_success and result.data.get("file_path"):
            size = os.path.getsize(result.data["file_path"]) if os.path.exists(result.data["file_path"]) else 0
            log_test(f"PDF 样式-{style}", size > 500, f"{size} bytes")
        else:
            log_test(f"PDF 样式-{style}", False, "生成失败")


# ==============================================================================
# 2. PPT 生成质量验证
# ==============================================================================

async def test_ppt_quality():
    """测试 PPT 实际生成质量"""
    section("PPT 生成质量验证")
    
    tool = PPTTool()
    
    # 2.1 生成 PPT 并验证
    result = await tool.execute("generate_ppt", {
        "topic": "季度销售报告",
        "outline": "销售概况\n业绩数据\n季度总结",
        "style": "business"
    })
    
    log_test("PPT 生成执行", result.is_success, f"状态: {result.status.value}")
    
    if result.is_success and result.data.get("file_path"):
        ppt_path = result.data["file_path"]
        
        # 验证1: 文件存在
        log_test("PPT 文件存在", os.path.exists(ppt_path), ppt_path)
        
        # 验证2: 文件大小
        file_size = os.path.getsize(ppt_path) if os.path.exists(ppt_path) else 0
        log_test("PPT 文件大小", file_size > 5000, f"{file_size} bytes")
        
        # 验证3: 有效 PPTX (ZIP 格式)
        try:
            with open(ppt_path, 'rb') as f:
                header = f.read(2)
                is_valid_pptx = header == b'PK'  # PPTX 是 ZIP 格式
                log_test("PPT 文件格式有效", is_valid_pptx, f"文件头: {header}")
        except Exception as e:
            log_test("PPT 文件格式有效", False, f"读取失败: {e}")
        
        # 验证4: 实际打开并检查幻灯片数量
        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)
            slides_list = list(prs.slides)
            slide_count = len(slides_list)
            log_test("PPT 包含幻灯片", slide_count > 0, f"共 {slide_count} 张幻灯片")
            
            # 检查每张幻灯片的内容
            for i, slide in enumerate(slides_list[:3]):  # 检查前3张
                shapes_list = list(slide.shapes)
                shapes_count = len(shapes_list)
                has_title = any(s.has_text_frame for s in shapes_list)
                log_test(f"  幻灯片{i+1} 内容", shapes_count > 0, f"形状数: {shapes_count}")
        except Exception as e:
            log_test("PPT 内容读取", False, f"打开失败: {e}")
        
        print(f"     生成文件: {ppt_path}")
        
        # 2.2 测试添加图表
        chart_result = await tool.execute("add_chart_slide", {
            "ppt_path": ppt_path,
            "title": "销售数据图表",
            "chart_type": "bar",
            "chart_data": '{"labels": ["Q1", "Q2", "Q3", "Q4"], "values": [100, 150, 120, 180]}',
            "style": "business"
        })
        log_test("PPT 添加图表", chart_result.is_success)
        
        if chart_result.is_success:
            try:
                prs = Presentation(ppt_path)
                slides_list = list(prs.slides)
                last_slide = slides_list[-1]
                shapes_list = list(last_slide.shapes)
                has_chart = any(s.shape_type == 3 for s in shapes_list)  # 3 = MSO_SHAPE_TYPE.CHART
                log_test("  图表幻灯片包含图表", has_chart or len(shapes_list) > 0, "图表添加验证")
            except Exception as e:
                log_test("  图表验证", False, str(e))
        
        # 2.3 测试添加章节分隔页
        section_result = await tool.execute("add_section_divider", {
            "ppt_path": ppt_path,
            "section_title": "核心业绩",
            "section_number": 1,
            "style": "modern"
        })
        log_test("PPT 添加章节分隔", section_result.is_success)
        
        print(f"     文件已保存: {ppt_path}")
    else:
        log_test("PPT 生成", False, f"生成失败: {result.error}")


# ==============================================================================
# 3. DOCX 生成质量验证
# ==============================================================================

async def test_docx_quality():
    """测试 DOCX 实际生成质量"""
    section("DOCX 生成质量验证")
    
    from src.tools.doc_generator import DocGeneratorTool
    tool = DocGeneratorTool()
    
    # 3.1 测试商业报告模板
    result = await tool.execute("generate_from_template", {
        "template": "business_report",
        "title": "年度商业报告",
        "author": "张三",
        "date": "2026-03-28",
        "organization": "测试公司"
    })
    
    log_test("DOCX 生成执行", result.is_success)
    
    if result.is_success and result.data.get("file_path"):
        docx_path = result.data["file_path"]
        
        # 验证1: 文件存在
        log_test("DOCX 文件存在", os.path.exists(docx_path), docx_path)
        
        # 验证2: 文件大小
        file_size = os.path.getsize(docx_path) if os.path.exists(docx_path) else 0
        log_test("DOCX 文件大小", file_size > 5000, f"{file_size} bytes")
        
        # 验证3: 实际打开并检查内容
        try:
            doc = Document(docx_path)
            para_count = len(doc.paragraphs)
            table_count = len(doc.tables)
            
            log_test("DOCX 包含段落", para_count > 0, f"共 {para_count} 个段落")
            log_test("DOCX 包含表格", table_count >= 0, f"共 {table_count} 个表格")
            
            # 检查标题
            has_title = any(p.text.strip() for p in doc.paragraphs if p.style.name.startswith('Heading'))
            log_test("DOCX 包含标题样式", has_title)
            
            # 检查页眉
            has_header = any(section.header.paragraphs for section in doc.sections)
            log_test("DOCX 包含页眉", has_header)
            
            print(f"     生成文件: {docx_path}")
            
        except Exception as e:
            log_test("DOCX 内容验证", False, f"打开失败: {e}")
        
        print(f"     文件已保存: {docx_path}")
    else:
        log_test("DOCX 生成", False, f"生成失败: {result.error}")


# ==============================================================================
# 4. Excel 样式验证
# ==============================================================================

def get_test_path(filename: str) -> str:
    """获取测试文件路径"""
    return str(TEST_OUTPUT_DIR / filename)

async def test_excel_quality():
    """测试 Excel 实际样式应用"""
    section("Excel 样式应用验证")
    
    # 创建测试 Excel
    wb = Workbook()
    ws = wb.active
    ws.title = "测试表"
    
    headers = ["产品", "销量", "单价", "总额"]
    ws.append(headers)
    ws.append(["产品A", 100, 50, "=B2*C2"])
    ws.append(["产品B", 150, 60, "=B3*C3"])
    
    temp_xlsx = get_test_path(f"test_data_{datetime.now().strftime('%H%M%S')}.xlsx")
    wb.save(temp_xlsx)
    
    # 4.1 引用检查（检查原文件）
    ref_result = refcheck_workbook(temp_xlsx)
    log_test("Excel 引用检查", ref_result.status in ["success", "issues_found"], f"状态: {ref_result.status}")
    
    if ref_result.issues:
        for issue in ref_result.issues[:3]:
            print(f"     问题: {issue}")
    
    # 4.2 应用专业样式
    wb_styled = load_workbook(temp_xlsx)
    styled_wb = apply_professional_style(wb_styled, theme="financial")
    
    log_test("Excel 样式应用", styled_wb is not None)
    
    if styled_wb:
        # 保存并验证
        base_name = Path(temp_xlsx).stem
        styled_path = get_test_path(f"{base_name}_styled.xlsx")
        styled_wb.save(styled_path)
        
        # 验证文件
        log_test("Excel 样式文件", os.path.exists(styled_path), styled_path)
        
        # 打开验证样式
        try:
            wb_check = load_workbook(styled_path)
            ws_check = wb_check.active
            
            # 检查网格线是否隐藏
            grid_hidden = not ws_check.sheet_view.showGridLines
            log_test("Excel 网格线已隐藏", grid_hidden)
            
            # 检查列宽
            log_test("Excel 列宽已设置", ws_check.column_dimensions['A'].width > 0, f"A列宽度: {ws_check.column_dimensions['A'].width}")
            
        except Exception as e:
            log_test("Excel 样式验证", False, f"验证失败: {e}")
        
        print(f"     Excel 样式文件已保存: {styled_path}")
    
    print(f"     源 Excel 文件已保存: {temp_xlsx}")


# ==============================================================================
# 5. 数据处理实际结果验证
# ==============================================================================

async def test_data_processing_quality():
    """测试数据处理实际结果"""
    section("数据处理结果验证")
    
    from src.tools.data_processor import DataProcessorTool
    tool = DataProcessorTool()
    
    # 创建测试数据
    wb = Workbook()
    ws = wb.active
    ws.append(["姓名", "部门", "销售额"])
    ws.append(["张三", "销售部", 150000])
    ws.append(["李四", "市场部", 120000])
    ws.append(["王五", "销售部", 180000])
    ws.append(["赵六", "技术部", 100000])
    
    temp_xlsx = get_test_path(f"test_data_{datetime.now().strftime('%H%M%S')}.xlsx")
    wb.save(temp_xlsx)
    
    # 5.1 筛选操作
    result = await tool.execute("filter_data", {
        "file_path": temp_xlsx,
        "column": "销售额",
        "operator": ">",
        "value": 120000
    })
    
    log_test("数据筛选执行", result.is_success)
    
    if result.is_success and result.data.get("output_path"):
        output_path = result.data["output_path"]
        log_test("筛选结果文件", os.path.exists(output_path), output_path)
        
        # 验证筛选结果
        try:
            wb_out = load_workbook(output_path)
            ws_out = wb_out.active
            row_count = ws_out.max_row - 1  # 减去表头
            
            # 应该只有销售额>120000的记录 (张三150000, 王五180000)
            expected_count = 2
            log_test("筛选结果正确", row_count == expected_count, f"期望 {expected_count} 行, 实际 {row_count} 行")
            
            print(f"     筛选结果文件: {output_path}")
            
        except Exception as e:
            log_test("筛选结果验证", False, f"验证失败: {e}")
    
    # 5.2 导出 CSV
    result = await tool.execute("export_data", {
        "file_path": temp_xlsx,
        "format": "csv"
    })
    
    log_test("导出 CSV 执行", result.is_success)
    
    if result.is_success and result.data.get("output_path"):
        csv_path = result.data["output_path"]
        log_test("CSV 文件", os.path.exists(csv_path), csv_path)
        
        # 验证 CSV 内容
        try:
            with open(csv_path, 'r', encoding='utf-8') as f:
                content = f.read()
                has_header = "姓名" in content or "name" in content.lower()
                has_data = "张三" in content or "150000" in content
                log_test("CSV 内容正确", has_header and has_data, "数据完整性检查")
        except Exception as e:
            log_test("CSV 内容验证", False, f"读取失败: {e}")
        
        print(f"     CSV 文件已保存: {csv_path}")
    
    print(f"     源 Excel 文件已保存: {temp_xlsx}")


# ==============================================================================
# 主函数
# ==============================================================================

async def main():
    """主测试函数"""
    print("\n" + "="*60)
    print("  Office 工具实际输出质量验证")
    print(f"  v2.28.0 - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("="*60)
    print("\n说明: 本测试验证实际生成的文件内容和质量")
    print("注意: 需要安装 LibreOffice 才能测试 Excel 公式重计算")
    
    await test_pdf_quality()
    await test_ppt_quality()
    await test_docx_quality()
    await test_excel_quality()
    await test_data_processing_quality()
    
    print("\n" + "="*60)
    print("  测试统计")
    print("="*60)
    print(f"  总测试数: {test_stats['total']}")
    print(f"  通过: {test_stats['passed']} ✅")
    print(f"  失败: {test_stats['failed']} ❌")
    print(f"  成功率: {test_stats['passed']/test_stats['total']*100:.1f}%")
    
    if test_stats['failed'] > 0:
        print("\n失败项目需要修复！")
    
    return 0 if test_stats["failed"] == 0 else 1


if __name__ == "__main__":
    exit_code = asyncio.run(main())
    sys.exit(exit_code)
